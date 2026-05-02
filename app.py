#!/usr/bin/env python3
import os
import re
import json
import zipfile
import subprocess
from io import BytesIO
import streamlit as st
from pathlib import Path
from datetime import datetime
import anthropic
from tavily import TavilyClient
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont

DESIGN_PRESETS = {
    "ダーク（黒背景・白文字）":   {"bg": "#1a1a1a", "title": "#ffffff", "content": "#e0e0e0", "accent": "#4a9eff"},
    "ホワイト（白背景・黒文字）": {"bg": "#ffffff", "title": "#1a1a1a", "content": "#333333", "accent": "#2563eb"},
    "ネイビー（紺背景・白文字）": {"bg": "#1e3a5f", "title": "#ffffff", "content": "#d0e8ff", "accent": "#ffd700"},
    "レッド（赤背景・白文字）":   {"bg": "#8b0000", "title": "#ffffff", "content": "#ffe0e0", "accent": "#ffcc00"},
    "グリーン（緑背景・白文字）": {"bg": "#1a4a2e", "title": "#ffffff", "content": "#d0ffd8", "accent": "#90ee90"},
}

SLIDE_FORMATS = {
    "横動画（16:9 / YouTube・一般動画）":  {"png": (1920, 1080), "pptx": (13.33, 7.5)},
    "正方形（1:1 / Instagram・note）":     {"png": (1080, 1080), "pptx": (10.0,  10.0)},
    "縦動画（9:16 / TikTok・Shorts・リール）": {"png": (1080, 1920), "pptx": (7.5,  13.33)},
}

SAMPLES_DIR = Path("samples")

PRESET_SIMPLE_KEYS = [
    "structure_type", "include_knowhow", "knowhow_theme", "knowhow_notes",
    "name", "category", "seller_name", "seller_first_person", "seller_profile",
    "interviewer_name", "interviewer_first_person", "interviewer_profile", "seller_authority",
    "seller_story", "mechanism",
    "catchcopy", "target_audience", "result1", "result2",
    "monthly_return", "ease_of_start",
    "voice1", "voice2", "voice3",
    "pain_points", "why_now", "vs_competition", "before_after",
    "third_party_type", "third_party_name", "third_party_points",
    "regular_price", "special_price", "limited_time", "limited_seats",
    "installment", "bonuses",
    "episode_structure", "closing_strength", "video_duration", "notes",
    "sales_flow_type", "sales_start_day", "consultation_method",
    "use_episode_themes",
]


def load_preset_to_session(preset):
    for k in PRESET_SIMPLE_KEYS:
        if k in preset:
            st.session_state[f"f_{k}"] = preset[k]
    for i, v in enumerate(preset.get("strengths", ["", "", "", ""])):
        st.session_state[f"str_{i}"] = v
    voices = preset.get("voices", ["", "", ""])
    for i in range(3):
        st.session_state[f"f_voice{i+1}"] = voices[i] if i < len(voices) else ""
    for i, v in enumerate(preset.get("comment_includes", [True]*5)):
        st.session_state[f"comment_include_{i}"] = v
    for i, v in enumerate(preset.get("comment_prompts", [""]*5)):
        st.session_state[f"comment_{i}"] = v
    for i, v in enumerate(preset.get("episode_themes", [""]*5)):
        st.session_state[f"episode_theme_{i}"] = v if i < 5 else ""


def clear_form_state():
    """フォームの入力欄をデフォルト値にリセットする。保存済みプリセットは残る。

    st.form 内のウィジェットはキーを削除しても表示がリセットされないため、
    load_preset_to_session と同様にデフォルト値を明示的にセットする方式を使う。
    """
    # テキスト入力・テキストエリア → 空文字
    for k in [
        "knowhow_theme", "knowhow_notes",
        "name", "seller_name", "seller_profile",
        "interviewer_name", "interviewer_profile", "seller_authority",
        "catchcopy", "target_audience",
        "result1", "result2", "monthly_return", "ease_of_start",
        "voice1", "voice2", "voice3",
        "pain_points", "why_now",
        "third_party_name", "third_party_points",
        "consultation_method",
        "regular_price", "special_price",
        "limited_time", "limited_seats", "bonuses", "notes",
        "seller_story", "mechanism", "vs_competition", "before_after",
    ]:
        st.session_state[f"f_{k}"] = ""

    # 強み欄（str_0〜str_3）
    for i in range(4):
        st.session_state[f"str_{i}"] = ""

    # コメント促進パート
    for i in range(5):
        st.session_state[f"comment_{i}"] = ""
        st.session_state[f"comment_include_{i}"] = True

    # 各話テーマ
    for i in range(5):
        st.session_state[f"episode_theme_{i}"] = ""

    # チェックボックス
    st.session_state["f_include_knowhow"]    = False
    st.session_state["f_use_episode_themes"] = False

    # ラジオボタン
    st.session_state["f_structure_type"]  = "従来型"
    st.session_state["f_sales_flow_type"] = "直接販売型"

    # セレクトボックス（フォームの default と一致させる）
    st.session_state["f_category"]          = "FX・為替投資"
    st.session_state["f_third_party_type"]  = "なし"
    st.session_state["f_sales_start_day"]        = "翌日（1日後）"
    st.session_state["f_installment"]            = "なし"
    st.session_state["f_episode_structure"]      = "1話完結"
    st.session_state["f_closing_strength"]       = "真摯・控えめ（押し付けない）"
    st.session_state["f_video_duration"]         = "7分（約2,100文字）"
    st.session_state["f_seller_first_person"]    = "私"
    st.session_state["f_interviewer_first_person"] = "私"


OUTPUT_DIR   = Path("output")
PRESETS_FILE = Path("presets.json")
MODEL = "claude-sonnet-4-6"

# ── 成功台本の設計図（心理構造フレームワーク） ──────────────────────────────
SUCCESS_FRAMEWORK = """
## 成功するプロダクトローンチ台本の構成要素（心理設計の型）

台本を書く際は、以下の8要素と話数設計を必ず意識して構成してください。

### 8つの必須構成要素

1. **衝撃的なベネフィットの提示（フック）**
   - 冒頭で「元手1万円で月60万」「30日で資金2倍」など常識外れの数字を提示する
   - 視聴者が「え、本当に？」と思わず続きを見てしまうフックを作る

2. **現状の問題提起と危機の共有**
   - 「普通に働いても生活を守れない時代」「年金制度の破綻」「AIによる職の代替」など
   - 視聴者が漠然と抱える不安を言語化し、共通の危機感を顕在化させる

3. **ストーリー（挫折→運命的な発見→成功）**
   - 発信者自身の「どん底（借金・大損・失職など）」を正直に告白する
   - そこから這い上がった「運命的な発見」を語り、親近感と信頼を同時に構築する

4. **非常識なロジック（独自メカニズム）の提示**
   - 「97%の負け組の逆を突く」「AIとプロの嗅覚の融合」など
   - 他と明確に違う"勝ち筋"を説明し、「これなら自分でもできる」という希望を与える

5. **圧倒的な証拠（エビデンス）の提示**
   - 取引履歴・LINEのやり取り・教え子の成功事例を動画や画像で見せる
   - 「見せられる証拠がある」ことで疑念を払拭する

6. **ハードルの徹底排除**
   - 「2万円から開始可能」「スマホ1台で完結」「設定代行サポートあり」など
   - 視聴者が「自分には無理」と思う理由をすべて先回りして潰す

7. **価格崩しと価値の正当化**
   - 本来100万円以上の価値があると印象づけた上で、補助金・期間限定などの理由で
   - 手が届く価格（9.8万〜30万）まで下げる演出をする

8. **強力なクロージング（希少性×緊急性）**
   - 「先着100名」「3日間限定」「審査制」などの枠を設け、今すぐ動く理由を作る
   - 5年後の未来を想像させ、行動しない損失を意識させる

---

### 話数別・最適設計図

#### 第1話：衝撃とパラダイムシフト
- オープニング：圧倒的な実績を見せ「あなたの常識を覆す」と宣言
- 質問フック：「月60万稼ぐのに元手はいくら必要だと思いますか？」など思考を揺さぶる
- 独自システムの紹介：少額からの可能性を示す
- 社会背景の危機：「沈みゆく日本」という共通の敵を設定し投資の必要性を説く
- 自己紹介（挫折ストーリー）：負けの過去を共有して信頼を獲得
- 次回予告＋コメント促進：「勝てる理由の全貌を次回明かす」として引きを作る

#### 第2話：信頼構築とメカニズム解明
- 反響の共有：「コメントが殺到している」として社会的証明を演出
- ロジックの深掘り：「市場の歪みを拾う仕組み」「AI24時間監視」など勝てる理屈を解説
- 第三者の登場：開発者・専門家・実績ある教え子との対談で客観性を担保
- 不安の先回り解消（Q&A形式）：初心者でも可能な理由、税金、リスクを事前に回答
- 特典の予告：「自己資金がない方向けの秘策」などを次話の引きとして提示

#### 第3話：オファーと決断の促し
- フルパッケージの公開：メインシステム＋長期・短期・複利など人生設計全体を提示
- 豪華特典の提示：導入マニュアル・無期限サポート・限定セミナーなどで価値を最大化
- 価格提示：他商品との比較・開発費の大きさを語った後、限定価格を発表
- リスクリワードの強調：「参加費は1週間で回収できる」という実績を再提示
- クロージング：5年後の未来を想像させ、先行受付URLへ誘導

#### 特別編（ダウンセル／救済クロージング）
- ハードル再調整：「価格で諦めた方へ」として機能を絞ったライト版・分割払いを再提案
- 取りこぼしの回収：本編で決断できなかった層に最後のチャンスを提供

---

### ライティング品質の基準

#### ① 数字の具体化ルール（必須）
抽象的な表現は必ず具体的な数字・期間・固有名詞に変換すること。

| NG（抽象）               | OK（具体）                        |
|--------------------------|-----------------------------------|
| 多くの人が               | 3,000名以上が                     |
| すぐに成果が出る         | 最短3日で、導入1週間で            |
| 大きな利益が出た         | 3ヶ月で31万円の利益               |
| 安く始められる           | 証拠金2万円から                   |
| 長年の実績               | 5年間・1,500本以上の動画で実証    |
| 高い勝率                 | 勝率78%、直近3ヶ月の実績          |

- 実績数値は「期間＋金額」をセットで使う
- 割合は「分母」も示す（「10人中8人」「導入者の85%が」）
- 曖昧な時制はNG → 「〇月〇日までの限定価格」

#### ② 感情の起伏設計（話数ごと）
視聴者の感情は「平坦」ではなく、意図的に上下させること。

- **第1話の感情グラフ**
  驚き・衝撃 → 危機感・焦り → 希望・ワクワク → 期待感（次回予告）

- **第2話の感情グラフ**
  共感・安心 → 知的納得 → 「自分でもできる」という確信 → 欲求の高まり

- **第3話の感情グラフ**
  価値への興奮 → 「今すぐ手に入れたい」 → 「やらない自分への後悔」 → 決断・行動

- **各話の感情ピーク**: 必ず1〜2箇所、視聴者が「おお！」と声に出すような強烈な一言・数字・ストーリーを配置すること

#### ③ 文章リズムの設計
台本は「読まれる文章」ではなく「聞かれる音声」。リズムが命。

- **基本パターン**: 短文 → 中文 → 短文 → 疑問文 → 短文
  例：「信じられないですよね。でも、これは実際に起きた話です。元手は、たった2万円でした。なぜ、こんなことが可能なのか？答えは、AIにあります。」

- **強調したい言葉の前後に必ず「間（ま）」を入れる**
  例：「この数字を見てください。……3ヶ月で、31万円。」

- **疑問文で視聴者を引き込む**
  セクションの冒頭は「〜だと思いませんか？」「〜したことはありますか？」で始めると視聴者が答えを探して聞き続ける

- **文末のバリエーション**: 「〜です」が3文連続するのはNG。「〜なんです」「〜ました」「〜ですよね」「〜でしょうか」を混ぜる
"""


def load_presets_from_file():
    if PRESETS_FILE.exists():
        try:
            return json.loads(PRESETS_FILE.read_text(encoding="utf-8"))
        except Exception:
            return {}
    return {}


def save_presets_to_file(presets):
    try:
        PRESETS_FILE.write_text(json.dumps(presets, ensure_ascii=False, indent=2), encoding="utf-8")
    except Exception:
        pass


def load_samples():
    SAMPLES_DIR.mkdir(exist_ok=True)
    samples = []
    # .txt優先：同名の.rtf/.rtfdがあっても.txtを使う
    txt_stems = {p.stem for p in SAMPLES_DIR.glob("*.txt")}
    files = list(SAMPLES_DIR.glob("*.txt"))
    for pattern in ["*.rtf", "*.rtfd"]:
        for p in SAMPLES_DIR.glob(pattern):
            if p.stem not in txt_stems:
                files.append(p)
    for file_path in sorted(set(files)):
        try:
            if file_path.suffix.lower() in (".rtf", ".rtfd"):
                result = subprocess.run(
                    ["textutil", "-convert", "txt", "-stdout", str(file_path)],
                    capture_output=True, text=True
                )
                content = result.stdout
            else:
                content = file_path.read_text(encoding="utf-8")
            if content.strip():
                samples.append({"filename": file_path.name, "content": content})
        except Exception:
            pass
    return samples


def build_system_prompt(samples):
    instruction_block = {
        "type": "text",
        "text": (
            "あなたは世界トップクラスのダイレクトレスポンス・コピーライター兼、敏腕映像ディレクターです。\n"
            "プロダクトローンチ用の動画台本を生成する専門家として、以下の設計図・サンプル・商品情報を参照してください。\n\n"
            "## 参照の優先度（必ず守ること）\n"
            "1. **サンプル台本の骨子・流れ**【最優先】：サンプルのセクション構成・話の展開順序・ナレーションの流れを「骨格」として最優先で使う。\n"
            "   サンプルがどこで何を語り、どう感情を動かし、どうCTAに向かっているか——その骨格を忠実に再現すること。\n"
            "2. **商品情報**：骨格の各セクションに、今回の商品固有の数字・ストーリー・特徴を肉付けする\n"
            "3. **成功台本の設計図**【補助のみ】：各セクションを「深掘り・強化」するための参考として使う。構成の骨格には使わない\n\n"
            "## 絶対に守るライティングルール\n"
            "- **書き言葉（〜である・〜のである・〜といえる・〜とされている）は絶対禁止**\n"
            "- 必ず「話し言葉」で書く：〜ですよね！ / 〜なんです / 〜じゃないですか / 〜してみてください\n"
            "- 1文はシンプルかつ短く。読点で区切り、一文に複数の意味を詰め込まない\n"
            "- 語尾は「〜ですよ。」を避け、「〜です。」「〜でした。」「〜ます。」「〜ました。」を基本にする\n"
            "- **人間味・自然さを絶対最優先する**：「台本を読んでいる感」が出た時点で失敗。\n"
            "  まるでその場で言葉を選びながら話しているように書くこと。言い直し・間・感情の揺れを入れる。\n"
            "  （使える表現例：「えっとですね…」「ちょっと待って、これ大事なんで」「あ、そうそう言い忘れてたんですけど」\n"
            "  「正直に言うと」「これ、実は私も最初は信じられなくて」「うまく言えないんですけど」）\n"
            "- **販売者のキャラクターを一貫して描く**：台本全体を通じて、その人固有の話し方・口癖・ユーモア・\n"
            "  弱さ・人間らしさを描くこと。誰が話しても同じになる「汎用コピー」にしないこと。\n"
            "- **温度感と共感を文章に乗せる**：視聴者の気持ちに寄り添う言葉を自然なタイミングで挟む。\n"
            "  （例：「わかります、私も最初はそう思ってました」「そうですよね、不安ですよね」「実は私も同じ悩みでした」）\n"
            "  一方的な説明にならないよう、視聴者との対話感を出すこと。\n"
            "- PREPの展開順序はあくまで深掘りの意識として使い、会話の流れを壊してまで型に当てはめないこと。\n"
            "- 動画冒頭の**開始15秒以内**に「この動画を最後まで見ないと損をする理由」を明示するフックを必ず入れること\n"
            "  （例：「この15秒だけ聞いてください。〇〇を知らないまま投資すると、9割の人が損します」）\n\n"
            "## 生成の進め方\n"
            "1. **サンプルの骨格を読み取る**：サンプルのセクション数・順序・各パートの役割・感情の流れを把握する\n"
            "2. **骨格に商品情報を当てはめる**：サンプルの各セクションの「場所」に、今回の数字・エピソード・強みを入れ替える\n"
            "3. **設計図で深掘りする**：設計図の8要素（フック・ストーリー・証拠・反論処理など）を使って各セクションを肉付けする\n"
            "- サンプルのトーン・語彙・言い回し・文章のリズムを踏襲する\n"
            "- 視聴者の心理変化（無関心→興味→信頼→欲求→行動）は、サンプルの流れに沿って自然に設計する\n"
        ),
    }

    framework_block = {
        "type": "text",
        "text": SUCCESS_FRAMEWORK,
        "cache_control": {"type": "ephemeral"},
    }

    if not samples:
        instruction_block["text"] += "\nサンプルがないため、設計図のフレームワークとベストプラクティスに従って生成してください。"
        return [instruction_block, framework_block]

    samples_text = "\n\n".join(
        f"=== サンプル {i + 1}: {s['filename']} ===\n{s['content']}"
        for i, s in enumerate(samples)
    )
    samples_block = {
        "type": "text",
        "text": f"## 学習用サンプル台本\n\n{samples_text}",
        "cache_control": {"type": "ephemeral"},
    }
    return [instruction_block, framework_block, samples_block]


def build_user_prompt(info, output_format="完全版"):
    strengths = [s for s in info.get("strengths", []) if s]
    strengths_str = "\n".join(f"  - {s}" for s in strengths) or "  （未入力）"

    voices = [v for v in info.get("voices", []) if v]
    voices_str = "\n".join(f"  - {v}" for v in voices) if voices else "  （なし）"

    structure_type = info.get("structure_type", "従来型")
    if structure_type == "従来型":
        structure_desc = "実績提示→商品説明→販売。インタビュー or 一人語り形式。"
    else:
        structure_desc = "無料ノウハウ・セミナー形式で価値提供→「時間・リスクが心配な方はこちら」と商品販売に自然につなぐフロントエンド型。"

    interviewer = info.get("interviewer_name", "").strip()
    dialogue_style = f"インタビュアー「{interviewer}」との対話形式" if interviewer else "一人語り（モノローグ）形式"

    comment_includes = info.get("comment_includes", [True] * 5)
    comment_prompts = info.get("comment_prompts", [])
    episode_structure = info.get("episode_structure", "1話完結")
    episode_count = int(episode_structure[0]) if episode_structure[0].isdigit() else 1

    include_knowhow = info.get("include_knowhow", False)
    knowhow_theme = info.get("knowhow_theme", "").strip()
    knowhow_notes = info.get("knowhow_notes", "").strip()

    lines = [
        "以下のプロダクト情報をもとに、**サンプル台本の骨子・セクション構成・話の展開を骨格として最優先に**台本を生成してください。",
        "サンプルの各セクションの順番・役割・感情の流れをそのまま骨格として使い、そこに以下の商品情報を肉付けして深掘りしてください。",
        "成功台本の設計図は骨格ではなく、各セクションを深掘り・強化するための補助ツールとして使ってください。",
        "",
        "## 構成タイプ",
        f"**タイプ**: {structure_type}",
        f"**説明**: {structure_desc}",
        f"**話し方スタイル**: {dialogue_style}",
        "",
        "## フロントエンド型ノウハウパート",
    ]

    if include_knowhow and knowhow_theme:
        lines += [
            f"**ノウハウテーマ**: {knowhow_theme}",
            f"**補足メモ**: {knowhow_notes if knowhow_notes else 'なし'}",
            "このテーマについて視聴者にとって価値ある具体的なノウハウを台本の前半に組み込んでください。",
            "ノウハウ提供の後、自然に商品の紹介につなげてください。",
        ]
    elif include_knowhow:
        lines.append("商品カテゴリに合った価値あるノウハウを台本の前半に自動生成して組み込んでください。")
    else:
        lines.append("ノウハウパートは不要です。")

    lines += [
        "",
        "## 商品・基本情報",
        f"**商品名**: {info.get('name', '')}",
        f"**ジャンル**: {info.get('category', '')}",
        f"**販売者名**: {info.get('seller_name', '')}",
        f"**販売者の一人称**: {info.get('seller_first_person', '私')}（台本全体でこの一人称に統一すること）",
        f"**販売者のプロフィール**: {info.get('seller_profile', '')}",
        f"**インタビュアーの一人称**: {info.get('interviewer_first_person', '私')}（インタビュアーの発言はこの一人称に統一すること）",
        f"**インタビュアープロフィール**: {info.get('interviewer_profile', '')}",
        f"**販売者の権威・実績**: {info.get('seller_authority', '')}",
        f"**原体験・どん底ストーリー**: {info.get('seller_story', '')}",
        f"**キャッチコピー**: {info.get('catchcopy', '')}",
        f"**ターゲット層**: {info.get('target_audience', '')}",
        f"**実績数値①（短期）**: {info.get('result1', '')}",
        f"**実績数値②（中〜長期）**: {info.get('result2', '')}",
        f"**月利 / 月収目安**: {info.get('monthly_return', '')}",
        f"**始めやすさの根拠**: {info.get('ease_of_start', '')}",
        f"**独自メカニズム・勝てる理由**: {info.get('mechanism', '')}",
        "",
        "**商品の強み**:",
        strengths_str,
        "",
        "## 利用者の声（喜びの声）",
        voices_str,
        "",
        "## 社会背景・痛み訴求",
        f"**視聴者のペイン**: {info.get('pain_points', '')}",
        f"**なぜ今必要か（why now）**: {info.get('why_now', '')}",
        f"**競合との違い・失敗体験との対比**: {info.get('vs_competition', '')}",
        f"**購入後のビフォーアフター**: {info.get('before_after', '')}",
        "",
        "## 第三者・信頼性",
        f"**第三者の種類**: {info.get('third_party_type', 'なし')}",
        f"**名前・肩書き**: {info.get('third_party_name', '')}",
        f"**裏付けポイント**: {info.get('third_party_points', '')}",
        "",
        "## 販売フロー",
    ]

    sales_flow = info.get("sales_flow_type", "直接販売型")
    lines.append(f"**販売フロータイプ**: {sales_flow}")
    if sales_flow == "直接販売型":
        start_day = info.get("sales_start_day", "翌日（1日後）")
        lines.append(f"**販売スタート日**: {start_day}")
        lines.append(f"クロージングでは、{start_day}から販売を開始することを明示し、視聴者に行動を促す緊迫感を持たせてください。")
    else:
        method = info.get("consultation_method", "").strip()
        if method:
            lines.append(f"**誘導先**: {method}")
        lines.append("クロージングでは商品への直接誘導ではなく、個別面談・相談への申し込みを促してください。「まずは無料で話しましょう」「個別でご相談ください」のようなトーンで、個別接触後に販売につなげる流れで台本を構成してください。")

    lines += [
        "",
        "## 価格・オファー",
        f"**定価**: {info.get('regular_price', '')}",
        f"**特別価格**: {info.get('special_price', '')}",
        f"**期間限定条件**: {info.get('limited_time', '')}",
        f"**先着・定員制限**: {info.get('limited_seats', '')}",
        f"**分割対応**: {info.get('installment', 'なし')}",
        f"**特典内容**: {info.get('bonuses', '')}",
        "",
        "## コメント促進パート（動画末尾）",
    ]

    any_comment = any(comment_includes[i] for i in range(episode_count))
    if not any_comment:
        lines.append("コメント促進パートは全話不要です。含めないでください。")
    else:
        for i in range(episode_count):
            include = comment_includes[i] if i < len(comment_includes) else True
            if not include:
                lines.append(f"**第{i+1}話のコメント促進**: 不要（含めないでください）")
            else:
                cp = comment_prompts[i].strip() if i < len(comment_prompts) else ""
                instruction = f"「{cp}」" if cp else "内容に合った質問を自動生成"
                lines.append(f"**第{i+1}話のコメント促進**: {instruction}")

    lines += [
        "",
        "## 構成オプション",
        f"**話数構成**: {episode_structure}",
        f"**1話あたりの動画の長さ**: {info.get('video_duration', '7分（約2,100文字）')}",
        f"**クロージングの強度**: {info.get('closing_strength', '標準')}",
    ]

    if info.get("use_episode_themes"):
        themes = info.get("episode_themes", [])
        specified = [(i+1, t) for i, t in enumerate(themes[:episode_count]) if t.strip()]
        if specified:
            lines.append("")
            lines.append("## 各話の内容指定")
            for num, theme in specified:
                lines.append(f"**第{num}話**: {theme}")
            lines.append("※ 内容が指定されている話は、その内容に沿って台本を作成してください。")

    lines += [
    ]

    if info.get("notes"):
        lines += ["", f"**追加メモ**: {info['notes']}"]

    seller_name_val = info.get("seller_name", "").strip()
    seller_pronoun = info.get("seller_first_person", "私")
    if seller_name_val:
        lines += [
            "",
            "## 自己紹介・一人称の厳守ルール",
            f"販売者の自己紹介シーンでは必ず「{seller_pronoun}は{seller_name_val}と申します」または「{seller_name_val}です」と名前を名乗ること。",
            f"台本全体で販売者の一人称は「{seller_pronoun}」に完全統一すること。他の一人称（私・僕・おれ・自分など）は一切使わないこと。",
        ]

    lines += [
        "",
        "【出力形式の指示】",
        "- **書き言葉（〜である・〜のである）は絶対禁止**。必ず話し言葉（〜ですよね・〜なんです）で書くこと",
        "- **抽象表現の完全禁止**：「素晴らしいです」「重要です」「人生が変わります」「大変でした」は一切使わない。視聴者が脳内で映像化できる具体描写で語ること",
        "- **深掘りの意識**：重要なトピックは【主張→理由→具体例・例え話→ベネフィット】の流れを意識して深く語ること。ただし会話の自然な流れを壊してまで型に当てはめないこと",
        *([
            "- **セリフのみをプレーンテキストで出力する**（テーブル不要）",
            "- 映像指示・演技指示は書かない。演者のセリフだけを自然な段落で書くこと",
            "- 話者が複数いる場合は行頭に「【販売者】」「【インタビュアー】」のように記載する",
        ] if output_format == "台本のみ" else [
            "- 台本は必ず以下の**3カラム Markdownテーブル形式**で出力すること：",
            "  | 映像/テロップ指示 | 演者のセリフ | 演技/効果音の指示 |",
            "  |---|---|---|",
            "  | [Bロール：〇〇] | 「〜ですよね。」 | [カメラ目線・ゆっくり] |",
            "  - 映像/テロップ: [Bロール：〇〇] [テロップ：〇〇] [画面切り替え] など",
            "  - 演技指示: [1秒の間] [カメラ目線で強調] [SE：〇〇] [場面転換] など",
        ]),
        "- 各セクションの前に ## セクション名（〜分〜秒） の見出しを入れること",
        "- 話数構成が指定されている場合は、その構成に合わせて台本を分けてください",
        "- コメント促進パートが必要な場合は各話の動画末尾に必ず含めてください",
        f"- 1話あたりの目標文字数を厳守してください（{info.get('video_duration', '7分（約2,100文字）')}）",
    ]
    return "\n".join(lines)


def build_info_from_session():
    """セッションステートからフォーム入力値を収集してinfoディクショナリを返す（フォーム外から呼ぶ用）。"""
    ss = st.session_state
    return {
        "structure_type":    ss.get("f_structure_type", "従来型"),
        "include_knowhow":   ss.get("f_include_knowhow", False),
        "knowhow_theme":     ss.get("f_knowhow_theme", ""),
        "knowhow_notes":     ss.get("f_knowhow_notes", ""),
        "name":              ss.get("f_name", ""),
        "category":          ss.get("f_category", ""),
        "seller_name":          ss.get("f_seller_name", ""),
        "seller_first_person":  ss.get("f_seller_first_person", "私"),
        "seller_profile":       ss.get("f_seller_profile", ""),
        "interviewer_name":     ss.get("f_interviewer_name", ""),
        "interviewer_first_person": ss.get("f_interviewer_first_person", "私"),
        "interviewer_profile":  ss.get("f_interviewer_profile", ""),
        "seller_authority":  ss.get("f_seller_authority", ""),
        "seller_story":      ss.get("f_seller_story", ""),
        "catchcopy":         ss.get("f_catchcopy", ""),
        "target_audience":   ss.get("f_target_audience", ""),
        "result1":           ss.get("f_result1", ""),
        "result2":           ss.get("f_result2", ""),
        "monthly_return":    ss.get("f_monthly_return", ""),
        "ease_of_start":     ss.get("f_ease_of_start", ""),
        "mechanism":         ss.get("f_mechanism", ""),
        "strengths":         [ss.get(f"str_{i}", "") for i in range(4)],
        "voices":            [ss.get(f"f_voice{i+1}", "") for i in range(3)],
        "pain_points":       ss.get("f_pain_points", ""),
        "why_now":           ss.get("f_why_now", ""),
        "vs_competition":    ss.get("f_vs_competition", ""),
        "before_after":      ss.get("f_before_after", ""),
        "third_party_type":  ss.get("f_third_party_type", "なし"),
        "third_party_name":  ss.get("f_third_party_name", ""),
        "third_party_points": ss.get("f_third_party_points", ""),
        "regular_price":     ss.get("f_regular_price", ""),
        "special_price":     ss.get("f_special_price", ""),
        "limited_time":      ss.get("f_limited_time", ""),
        "limited_seats":     ss.get("f_limited_seats", ""),
        "installment":       ss.get("f_installment", "なし"),
        "bonuses":           ss.get("f_bonuses", ""),
        "comment_includes":  [ss.get(f"comment_include_{i}", True) for i in range(5)],
        "comment_prompts":   [ss.get(f"comment_{i}", "") for i in range(5)],
        "episode_structure": ss.get("f_episode_structure", "1話完結"),
        "closing_strength":  ss.get("f_closing_strength", "標準"),
        "video_duration":    ss.get("f_video_duration", "7分（約2,100文字）"),
        "use_episode_themes": ss.get("f_use_episode_themes", False),
        "episode_themes":    [ss.get(f"episode_theme_{i}", "") for i in range(5)],
        "notes":             ss.get("f_notes", ""),
        "sales_flow_type":   ss.get("f_sales_flow_type", "直接販売型"),
        "sales_start_day":   ss.get("f_sales_start_day", "翌日（1日後）"),
        "consultation_method": ss.get("f_consultation_method", ""),
    }


def build_outline_prompt(info):
    """構成案（アウトライン）生成用プロンプト。台本の文章は書かせない。"""
    base = build_user_prompt(info)
    episode_structure = info.get("episode_structure", "1話完結")
    episode_count = int(episode_structure[0]) if episode_structure[0].isdigit() else 1
    video_duration = info.get("video_duration", "7分（約2,100文字）")

    structure_type = info.get("structure_type", "従来型")
    is_frontend = structure_type == "フロントエンド型"

    if is_frontend and episode_count >= 2:
        # フロントエンド型専用の話数別ルール
        _fe_rules = {
            2: """
### 2話構成（フロントエンド型）の厳守ルール（違反禁止）
- **第1話（ノウハウ・価値提供）**: 商品名・価格・売り込みは一切禁止。視聴者にとって価値あるノウハウを惜しみなく提供し、信頼を構築する。最後に「次回、さらに効率的に〇〇できる特別な方法を公開します」と期待させて終わること。
- **第2話（商品紹介・クロージング）**: 「ノウハウを自力でやり続けるのは大変→商品を使えばより楽に・早く・確実に結果が出る」という自然な流れで商品を提案。特典・保証・価格の正当性・今すぐ行動すべき理由でクロージングすること。""",
            3: """
### 3話構成（フロントエンド型）の厳守ルール（違反禁止）
- **第1話（問題提起・共感・ノウハウ①）**: 視聴者の痛みに深く共感し、問題の根本原因を解説。ノウハウの入口部分を提供して信頼を構築。商品・価格は禁止。「次回は〇〇の核心を公開します」と期待させて終わること。
- **第2話（ノウハウ②・深掘り）**: ノウハウの核心部分を惜しみなく提供。成功事例・失敗事例を交えて価値を最大化。商品・価格は禁止。「次回、これを最短で実現する方法を発表します」と強烈に予告して終わること。
- **第3話（商品紹介・クロージング）**: 「自力でやる大変さ→商品なら効率的」という流れで自然に商品提案。特典・保証・価格の正当性・緊急性でクロージングすること。""",
            4: """
### 4話構成（フロントエンド型）の厳守ルール（違反禁止）
- **第1話（問題提起・共感）**: 視聴者の痛みに深く共感し、なぜうまくいかないかの根本原因を解説。商品・価格は禁止。次回への期待で終わること。
- **第2話（ノウハウ①・解決策の方向性）**: 解決策の全体像とステップ①を具体的に解説。「自分でもできそう」と思わせる。商品・価格は禁止。
- **第3話（ノウハウ②・反論処理）**: ステップ②③の深掘りと、よくある疑問・反論を論理的に潰す。商品・価格は禁止。「次回、特別なサポートプログラムを発表します」と強烈に予告して終わること。
- **第4話（商品紹介・クロージング）**: 「自力でやる苦労→商品なら楽・早い・確実」という流れで自然に商品提案。特典・保証・価格・緊急性でクロージングすること。""",
            5: """
### 5話構成（フロントエンド型）の厳守ルール（違反禁止）
- **第1話（問題提起・共感）**: 視聴者の痛みに深く共感し、問題の根本原因を解説。商品・価格は禁止。次回への期待で終わること。
- **第2話（ノウハウ①）**: 解決策のステップ①を具体的に解説。成功事例を交えて信頼構築。商品・価格は禁止。
- **第3話（ノウハウ②）**: ステップ②③の深掘り。「自分でもできる」という自信を持たせる。商品・価格は禁止。
- **第4話（ノウハウ③・反論処理）**: ステップ④⑤と、よくある疑問・反論を潰す。商品・価格は禁止。「次回、特別なプログラムを発表します」と強烈に予告して終わること。
- **第5話（商品紹介・クロージング）**: 「自力でやる苦労→商品なら効率的」という流れで自然に商品提案。特典・保証・価格・緊急性でクロージングすること。""",
        }
        episode_rules = _fe_rules.get(episode_count, f"""
### {episode_count}話構成（フロントエンド型）の厳守ルール（違反禁止）
- **第1話〜第{episode_count-1}話（ノウハウ・価値提供）**: 各話で役立つノウハウを段階的に提供。商品名・価格・売り込みは一切禁止。各話の末尾で次回への期待を持たせること。
- **第{episode_count}話（商品紹介・クロージング）**: 「自力でやる大変さ→商品なら効率的」という流れで自然に商品提案。特典・保証・価格・緊急性でクロージングすること。""")
    elif episode_count == 3:
        episode_rules = """
### 3話構成の厳守ルール（違反禁止）
- **第1話（機会の提示）**: 価格・売り込みは絶対禁止。パラダイムシフトとストーリーに80%の時間を割く。「次回は具体的なステップを公開します」と期待させて終わること。
- **第2話（教育と反論処理）**: ノウハウ解説と視聴者の疑問を論理的に潰す反論処理を組み合わせる。最後に「次回、あなたを直接サポートする特別なプログラムを発表します」と強烈に予告して終わること。
- **第3話（セールス）**: 共感は手短にし、「商品内容」「圧倒的な特典」「保証」「価格の正当性」「今すぐ買うべき理由（希少性・期限）」に焦点を当てること。"""
    elif episode_count == 4:
        episode_rules = """
### 4話構成の厳守ルール（違反禁止）
- **第1話（機会の提示）**: 価格・売り込みは絶対禁止。パラダイムシフトとストーリーに80%の時間を割く。次回への期待で終わる。
- **第2話（教育とノウハウ）**: ノウハウ・独自メソッドの解説に重点。図解・事例を多用し売り込みはしない。
- **第3話（未来の提示と反論処理）**: 疑問を論理的に潰す反論処理を中心に構成。「次回、特別なプログラムを発表します」と強烈に予告。
- **第4話（セールス）**: 商品内容、特典、保証、価格の正当性、今すぐ買うべき理由に焦点を当てたセールス構成。"""
    else:
        episode_rules = ""

    return f"""{base}

---

## 【出力指示：構成案（アウトライン）のみを出力してください】

台本の文章は**絶対に書かないでください**。
以下の形式で「骨格（アウトライン）」だけを出力してください。

### 構成案の作り方（この順序で考えること）
1. **まずサンプル台本の骨格を踏襲する**：サンプルのセクション数・順序・各パートの役割・話の展開パターンをそのまま骨格として使うこと
2. **各セクションに商品情報を当てはめる**：骨格の「場所」に今回の商品の数字・エピソード・強みを入れ替えて深掘り素材を考える
3. **成功台本の設計図で肉付けする**：フック・ストーリー・反論処理などの要素で各セクションを強化する
{episode_rules}

各セクションについて記載すること：
1. セクション名・時間目安・**時間配分（%）**
2. このセクションで伝えるポイント（ロジックの展開：箇条書き3〜5個）
3. 使うべき具体的な数字・実績・エピソード
4. 【深掘り素材】このブロックを語るための以下のアイデアを必ず書き出すこと：
   - 使える「具体的なエピソード」（before/after・失敗談・顧客事例など）
   - 使える「例え話・比喩」（〜に例えると／まるで〜のように）
   - 裏付けとなる「データ・事実・統計」
   ※ここで素材が薄ければSTEP2の台本も薄くなる。深掘りして捻出すること
5. 感情的な目標（視聴者にどう感じさせるか）

---

出力フォーマット：

# 第1話：「[タイトル案]」（{video_duration}）
**この話の役割**：機会の提示 / 教育 / セールス など

## [1] オープニング・衝撃フック（時間配分：約15%）
- 時間目安：〜XX分
- ポイント：〇〇
- 使う数字/エピソード：〇〇
- 感情目標：視聴者が「え、本当に？続きが見たい」と感じる

## [2] 社会背景・問題提起（時間配分：約20%）
...

---

{episode_count}話分の構成案を全て出力してください。
最後に「---構成案ここまで---」と書いてください。"""


def build_script_from_outline_prompt(outline, info):
    """承認済み構成案をもとに台本を生成するプロンプト。"""
    video_duration = info.get("video_duration", "7分（約2,100文字）")
    min_chars = min_chars_per_block(video_duration)
    closing_strength = info.get("closing_strength", "標準")
    interviewer = info.get("interviewer_name", "").strip()
    dialogue_style = f"インタビュアー「{interviewer}」との対話形式" if interviewer else "一人語り（モノローグ）形式"
    name = info.get("name", "")
    seller_name = info.get("seller_name", "")
    result1 = info.get("result1", "")
    result2 = info.get("result2", "")
    voices = [v for v in info.get("voices", []) if v]
    voices_str = "\n".join(f"  - {v}" for v in voices) if voices else "  （なし）"

    return f"""以下の【確定した構成案】に**厳密に従って**、プロダクトローンチ用の動画台本を生成してください。

## 確定した構成案（この構成・順序・ポイントを全て反映すること）
{outline}

---

## 商品情報（参照用）
- 商品名：{name}
- 販売者名：{seller_name}
- 実績①：{result1}
- 実績②：{result2}
- 利用者の声：
{voices_str}

---

## 生成ルール
- 構成案のセクション順・ポイント・数字をすべて台本に反映すること
- 1話あたりの目安：{video_duration}
- 話し方スタイル：{dialogue_style}
- クロージングの強度：{closing_strength}
- **書き言葉（〜である・〜のである）は絶対禁止**。必ず話し言葉で書くこと
- **抽象表現の完全禁止**：「素晴らしいです」「重要です」「人生が変わります」「大変でした」は一切使わない。視聴者が脳内で映像化できるレベルの具体描写で語ること
- **深掘りの意識**：重要なトピックは【主張→理由→具体的なエピソード・例え話→視聴者へのベネフィット】の流れを意識して深く語ること。ただし会話の自然な流れを最優先し、型に機械的に当てはめないこと
- **最低文字数**：1セクションにつき最低{min_chars}文字以上を目安に展開すること

## 出力フォーマット（3カラム Markdownテーブル形式）

台本は必ず以下の3カラム表形式で出力すること：

| 映像/テロップ指示 | 演者のセリフ | 演技/効果音の指示 |
|---|---|---|
| [Bロール：〇〇の様子] | 「〜ですよね。」 | [カメラ目線・ゆっくり話す] |

- **映像/テロップ指示**: `[Bロール：〇〇]` `[テロップ：〇〇]` `[画面切り替え]` `[静止画：〇〇]` など
- **演者のセリフ**: 話し言葉のみ。1行1文を基本とし、自然な間で改行する
- **演技/効果音の指示**: `[1秒の間]` `[カメラ目線で強調]` `[SE：ドラム音]` `[場面転換]` `[感情を込めて]` など

各セクションの前に `## セクション名（〜分〜秒）` の見出しを必ず入れること。
"""


def split_script_by_episode(script):
    """台本を # 第N話 ヘッダーで話ごとに分割する。分割できない場合は全体を1件で返す。"""
    lines = script.split('\n')
    episodes = []
    current_title = None
    current_lines = []
    for line in lines:
        if re.match(r'^#\s+第\d+話', line.strip()):
            if current_lines and any(l.strip() for l in current_lines):
                episodes.append({
                    'title': current_title or '全話',
                    'content': '\n'.join(current_lines).strip(),
                })
            current_title = line.strip().lstrip('#').strip()
            current_lines = [line]
        else:
            current_lines.append(line)
    if current_lines and any(l.strip() for l in current_lines):
        episodes.append({
            'title': current_title or '全話',
            'content': '\n'.join(current_lines).strip(),
        })
    return episodes if episodes else [{'title': '全話', 'content': script}]


def min_chars_per_block(video_duration):
    """動画の尺から1ブロックあたりの最低文字数を返す。"""
    for minutes, chars in [(120, 2500), (90, 2000), (60, 1500), (45, 1200),
                            (30, 900), (20, 700), (15, 600)]:
        if f"{minutes}分" in video_duration:
            return chars
    return 400


def parse_outline_blocks(outline_text):
    """構成案を ## [N] で始まるブロック単位に分割する。"""
    blocks = []
    current_title = ""
    current_lines = []
    episode_header = ""
    for line in outline_text.split('\n'):
        s = line.strip()
        if re.match(r'^# ', s) and '第' in s:
            episode_header = s
        elif re.match(r'^## \[', s):
            if current_title and current_lines:
                blocks.append({
                    'title': current_title,
                    'episode_header': episode_header,
                    'content': '\n'.join(current_lines),
                })
            current_title = s
            current_lines = [line]
        else:
            current_lines.append(line)
    if current_title and current_lines:
        blocks.append({
            'title': current_title,
            'episode_header': episode_header,
            'content': '\n'.join(current_lines),
        })
    return blocks


def build_block_script_prompt(outline, block, block_num, total_blocks, previous_script, info, output_format="完全版"):
    """1ブロック分の台本を生成するプロンプト。"""
    video_duration = info.get("video_duration", "7分（約2,100文字）")
    min_chars = min_chars_per_block(video_duration)
    interviewer = info.get("interviewer_name", "").strip()
    dialogue_style = f"インタビュアー「{interviewer}」との対話形式" if interviewer else "一人語り（モノローグ）形式"
    episode_header = block.get('episode_header', '')
    block_label = f"{episode_header} / {block['title']}" if episode_header else block['title']
    prev_context = ""
    if previous_script:
        tail = previous_script[-600:]
        prev_context = f"\n## 直前ブロックの末尾（流れを繋げるために参照）\n{tail}\n\n---\n"

    if output_format == "台本のみ":
        format_rule = (
            "- **セリフのみをプレーンテキストで出力する**（テーブル不要）\n"
            "- 映像指示・演技指示は書かない。演者のセリフだけを自然な段落で書くこと\n"
            "- 各発言の話者名が複数いる場合は「【販売者】」「【インタビュアー】」のように行頭に記載する"
        )
    else:
        format_rule = (
            "- **3カラム Markdownテーブル形式**で出力する：\n"
            "  | 映像/テロップ指示 | 演者のセリフ | 演技/効果音の指示 |\n"
            "  |---|---|---|"
        )

    return f"""以下の【全体構成案】の中から、**今回のブロックのみ**の台本を生成してください。他のブロックは書かないこと。

## 全体構成案（参照用）
{outline}

---
{prev_context}
## 今回書くブロック（{block_num} / {total_blocks}）
{block_label}

{block['content']}

---

## 生成ルール（厳守）
- **このブロックの台本のみ**を出力する（他のブロックは書かない）
- **最低{min_chars}文字以上**で重厚に展開すること（薄い内容は絶対禁止）
- **深掘りの意識**：重要なトピックは【主張→理由→具体的なエピソード・例え話→ベネフィット】の流れを意識して語ること。型に機械的に当てはめず、会話の自然な流れを最優先にすること
- **抽象表現は避ける**：「素晴らしいです」「重要です」「大変でした」より、視聴者が脳内で映像化できる具体描写を使うこと
- 話し方スタイル：{dialogue_style}
{format_rule}
- 書き言葉（〜である）絶対禁止、話し言葉のみ
"""


def try_parse_json(text):
    """テキストからJSONを抽出し、一般的な破損パターンを自動修復してパースする。"""
    match = re.search(r'\{.*\}', text, re.DOTALL)
    if not match:
        raise ValueError("JSONが見つかりませんでした")
    raw = match.group()

    # ① そのままパース
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        pass

    # ② よくある修正: trailing comma / 隣接オブジェクト間の欠落カンマ
    cleaned = re.sub(r',(\s*[}\]])', r'\1', raw)          # trailing comma
    cleaned = re.sub(r'}\s*\n(\s*)\{', r'},\n\1{', cleaned)  # } { → },{
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        pass

    # ③ 末尾トランケーション修復: 最後の完全スライドで切り詰めて閉じ括弧を補完
    for sep in ('}\n    }', '},\n    {', '},\n  {', '},\n{', '},'):
        pos = cleaned.rfind(sep)
        if pos > 0:
            truncated = cleaned[:pos + 1]
            open_b = truncated.count('[') - truncated.count(']')
            open_c = truncated.count('{') - truncated.count('}')
            repaired = truncated
            repaired += ']' * max(0, open_b)
            repaired += '}' * max(0, open_c)
            try:
                return json.loads(repaired)
            except json.JSONDecodeError:
                continue

    raise ValueError(
        "スライドJSONの解析に失敗しました。\n"
        "台本が非常に長い場合、スライド数が多すぎてJSONが壊れることがあります。\n"
        "話数構成を減らすか台本を短くしてお試しください。"
    )


SLIDE_CHUNK_MAX_CHARS = 4500  # このサイズ以下に収まるよう台本を分割（約15分相当）


def split_script_into_chunks(script):
    """台本を話数・長さで分割してチャンクリストを返す。

    戻り値: [(episode_num, chunk_label, chunk_text), ...]
    - 第N話マーカーがあれば話数ごとに分割
    - 1チャンクが SLIDE_CHUNK_MAX_CHARS を超える場合は前半/後半/パートN でさらに分割
    """
    # ── 話数マーカーで分割 ──
    ep_pattern = re.compile(r'(?:^|\n)(【?第\s*(\d+)\s*話[^\n]*)', re.MULTILINE)
    ep_matches = list(ep_pattern.finditer(script))

    if len(ep_matches) >= 2:
        ep_texts = {}
        for i, m in enumerate(ep_matches):
            ep_num = int(m.group(2))
            start = m.start()
            end = ep_matches[i + 1].start() if i + 1 < len(ep_matches) else len(script)
            ep_texts[ep_num] = script[start:end].strip()
    else:
        ep_texts = {1: script}

    # ── 各話を長さでさらに分割 ──
    result = []
    for ep_num, ep_text in sorted(ep_texts.items()):
        ep_label_base = f"第{ep_num}話" if len(ep_texts) > 1 else "全体"

        if len(ep_text) <= SLIDE_CHUNK_MAX_CHARS:
            result.append((ep_num, ep_label_base, ep_text))
            continue

        # 段落単位で SLIDE_CHUNK_MAX_CHARS 以下のチャンクに分割
        paras = re.split(r'\n\n+', ep_text)
        chunks_text, cur, cur_len = [], [], 0
        for para in paras:
            if cur_len + len(para) > SLIDE_CHUNK_MAX_CHARS and cur:
                chunks_text.append('\n\n'.join(cur))
                cur, cur_len = [para], len(para)
            else:
                cur.append(para)
                cur_len += len(para)
        if cur:
            chunks_text.append('\n\n'.join(cur))

        n = len(chunks_text)
        part_labels = (
            ["前半", "後半"] if n == 2
            else [f"パート{j+1}" for j in range(n)]
        )
        for ci, chunk in enumerate(chunks_text):
            label = f"{ep_label_base} {part_labels[ci]}"
            result.append((ep_num, label, chunk))

    return result


def search_trends(tavily_api_key, category, product_name):
    """カテゴリと商品名に関連する最新トレンドをWeb検索して返す。"""
    try:
        client = TavilyClient(api_key=tavily_api_key)
        query = f"{category} {product_name} 最新トレンド 2025 日本"
        results = client.search(query=query, max_results=3, search_depth="basic")
        summaries = []
        for r in results.get("results", []):
            title = r.get("title", "")
            content = r.get("content", "")[:300]
            summaries.append(f"・{title}：{content}")
        return "\n".join(summaries) if summaries else ""
    except Exception:
        return ""


def generate_slide_data(client, script):
    """Claudeに台本を渡してスライド構成をJSON形式で生成する（Adaptive Thinking使用）。"""
    char_count = len(script.replace(" ", "").replace("\n", ""))
    target_slides = max(15, char_count // 120)

    prompt = f"""以下のプロダクトローンチ台本を、動画に同期するプレゼンテーションスライドに変換してください。

## 台本
{script}

---

## スライド設計の原則

### 枚数・密度
- 目標: **{target_slides}枚以上**（台本の話題転換・セクション区切りごとに1枚）
- 1スライド1メッセージ。詰め込みすぎない

### レイアウト使い分け（必ず適切に選ぶこと）
- **impact**: 数字・価格・CTA・感情的クライマックス・衝撃的な一言。全体の約30%
- **section**: 大きな話題転換・章の区切り・「では次に〜」の瞬間。全体の約20%
- **standard**: 複数ポイントの説明・機能・利用者の声・Q&A。全体の約50%

### タイトルの書き方（最重要）
- **20文字以内の断言文**
- NG（抽象）: 「商品の特徴について」
- OK（具体）: 「元手2万円で始められる」「97%が知らない勝ち筋」
- 数字・疑問・感嘆を積極活用: 「月利10%、3ヶ月で実証」「なぜ今すぐ動くべきか」

### contentの書き方
- impact / section: 0〜1個（補足のみ、なくてもよい）
- standard: 2〜4個の箇条書き
- 各項目は**20文字以内**、体言止めか短い断言文
- NG: 「〜することができます」 → OK: 「〜を実現」「〜の秘密」「〜で解決」

### 絵文字
- 台本の感情・内容に正確に合ったもの1つ
- 数字スライドは📊💰、感情スライドは😤💡、CTAは🔥⚡ など

---

## 出力形式（JSONのみ。説明文・コードブロック記号不要）

{{
  "title": "動画シリーズのタイトル",
  "slides": [
    {{
      "episode": 1,
      "layout": "impact",
      "title": "スライドタイトル（20文字以内）",
      "content": ["補足テキスト（短く）"],
      "emoji": "📈",
      "notes": "台本の該当箇所の抜粋"
    }}
  ]
}}

`episode`は必須。そのスライドが何話目かを整数で（1話完結ならすべて1）。"""

    with client.messages.stream(
        model=MODEL,
        max_tokens=32000,
        thinking={"type": "adaptive"},
        messages=[{"role": "user", "content": prompt}]
    ) as stream:
        response = stream.get_final_message()
    # Adaptive Thinking使用時はthinkingブロックとtextブロックが混在する
    text = next((b.text for b in response.content if b.type == "text"), "")
    if not text:
        raise ValueError("スライドデータの生成に失敗しました")
    return try_parse_json(text)


def build_pptx(slide_data, design, pptx_size=(13.33, 7.5)):
    """高品質PPTXを生成する。レイアウト3種・装飾要素・テキスト階層を強化。"""
    W, H = pptx_size
    landscape = W >= H

    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    def _rgb(hex_color):
        h = hex_color.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    bg_rgb  = _rgb(design["bg"])
    tc_rgb  = _rgb(design["title"])
    cc_rgb  = _rgb(design["content"])
    acc_rgb = _rgb(design["accent"])

    def set_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_rgb

    def R(slide, l, t, w, h, color):
        """塗りつぶし矩形を追加する。"""
        shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        return shp

    def T(slide, text, l, t, w, h, pt, color, bold=False, align=PP_ALIGN.LEFT):
        """テキストボックスを追加する。"""
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(pt)
        run.font.color.rgb = color
        run.font.bold = bold
        return tb

    def B(slide, items, l, t, w, h, pt, color, align=PP_ALIGN.CENTER):
        """箇条書きテキストボックスを追加する（行間・段落間を整える）。"""
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_before = Pt(10)
            p.space_after  = Pt(2)
            run = p.add_run()
            run.text = item
            run.font.size = Pt(pt)
            run.font.color.rgb = color

    # ── サイズ適応フォント・寸法 ──
    if landscape:
        f_hero   = 56   # impactメインタイトル
        f_title  = 42   # section/standardタイトル
        f_body   = 26   # 本文・箇条書き
        f_sub    = 22   # サブテキスト
        bar      = 0.10 # 上部アクセントバー（薄め）
        sep      = 0.04 # 細いセパレータ
        pad      = 0.7  # 左右パディング
    else:
        f_hero   = 40
        f_title  = 30
        f_body   = 20
        f_sub    = 17
        bar      = 0.08
        sep      = 0.03
        pad      = 0.5

    # ════════════════════════════════════════
    # タイトルスライド（センター・大）
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    R(slide, 0, 0,          W, bar * 2, acc_rgb)              # 上バー
    R(slide, 0, H - bar * 2, W, bar * 2, acc_rgb)             # 下バー
    T(slide, slide_data.get("title", ""),
      pad, H * 0.18, W - pad * 2, H * 0.55,
      f_hero, tc_rgb, bold=True, align=PP_ALIGN.CENTER)
    R(slide, W * 0.30, H * 0.76, W * 0.40, sep * 1.5, acc_rgb)  # 装飾ライン

    # ════════════════════════════════════════
    # コンテンツスライド
    # ════════════════════════════════════════
    for slide_info in slide_data.get("slides", []):
        slide = prs.slides.add_slide(blank)
        set_bg(slide)

        layout  = slide_info.get("layout", "standard")
        emoji   = slide_info.get("emoji", "")
        title   = slide_info.get("title", "")
        content = slide_info.get("content", [])

        # 共通：上部アクセントバー
        R(slide, 0, 0, W, bar, acc_rgb)

        # ────────────────────────────────────
        # IMPACT レイアウト
        # 上下バー + 特大センタータイトル
        # ────────────────────────────────────
        if layout == "impact":
            R(slide, 0, H - bar * 2, W, bar * 2, acc_rgb)     # 下バー（太め）

            has_sub = bool(content)
            if emoji:
                T(slide, emoji,
                  pad, H * 0.16, W - pad * 2, bar * 3,
                  f_title, acc_rgb, align=PP_ALIGN.CENTER)
                ty = H * 0.30
            else:
                ty = H * 0.22 if has_sub else H * 0.26

            T(slide, title,
              pad, ty, W - pad * 2, H * 0.42,
              f_hero, tc_rgb, bold=True, align=PP_ALIGN.CENTER)
            R(slide, W * 0.35, ty + H * 0.44, W * 0.30, sep, acc_rgb)

            if has_sub:
                sub = "　•　".join(content[:2])
                T(slide, sub,
                  pad, ty + H * 0.47, W - pad * 2, 0.7,
                  f_sub, cc_rgb, align=PP_ALIGN.CENTER)

        # ────────────────────────────────────
        # SECTION レイアウト
        # 薄い上バー + 大センタータイトル + アクセントライン
        # ────────────────────────────────────
        elif layout == "section":
            if emoji:
                T(slide, emoji,
                  pad, H * 0.17, W - pad * 2, bar * 3.5,
                  f_title, acc_rgb, align=PP_ALIGN.CENTER)
                ty = H * 0.34
            else:
                ty = H * 0.28

            T(slide, title,
              pad, ty, W - pad * 2, H * 0.38,
              f_title + 6, tc_rgb, bold=True, align=PP_ALIGN.CENTER)
            R(slide, W * 0.30, ty + H * 0.40, W * 0.40, sep * 1.5, acc_rgb)

            if content:
                T(slide, content[0],
                  pad, ty + H * 0.44, W - pad * 2, 0.6,
                  f_sub, cc_rgb, align=PP_ALIGN.CENTER)

        # ────────────────────────────────────
        # STANDARD レイアウト
        # 薄い上バー + センタータイトル + 区切り線 + 箇条書き（中央）
        # ────────────────────────────────────
        else:
            title_str = f"{emoji}  {title}" if emoji else title
            T(slide, title_str,
              pad, bar + 0.18, W - pad * 2, 1.0,
              f_title, tc_rgb, bold=True, align=PP_ALIGN.CENTER)

            sep_y = bar + 1.25
            R(slide, pad, sep_y, W - pad * 2, sep, acc_rgb)

            if content:
                B(slide, content,
                  pad, sep_y + 0.20, W - pad * 2, H - sep_y - 0.30,
                  f_body, cc_rgb, align=PP_ALIGN.CENTER)

        # スピーカーノート
        notes = slide_info.get("notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def hex_to_rgb(hex_color):
    h = hex_color.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def _find_cjk_font():
    """Find a CJK-capable font on this system and cache the path."""
    candidates = [
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansJP-Regular.otf",
        "/usr/share/fonts/truetype/noto/NotoSansJP-Regular.ttf",
        "/System/Library/Fonts/ヒラギノ角ゴシック W3.ttc",
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
    ]
    for p in candidates:
        if os.path.exists(p):
            return p
    # Fallback: glob search for any Noto CJK font
    import glob
    for pattern in [
        "/usr/share/fonts/**/*CJK*Regular*.ttc",
        "/usr/share/fonts/**/*CJK*Regular*.otf",
        "/usr/share/fonts/**/*Noto*Regular*.ttc",
    ]:
        matches = sorted(glob.glob(pattern, recursive=True))
        if matches:
            return matches[0]
    return None


_FONT_PATH = _find_cjk_font()


def get_font(size):
    if _FONT_PATH:
        try:
            return ImageFont.truetype(_FONT_PATH, size)
        except Exception:
            pass
    return ImageFont.load_default()


def wrap_text(text, font, max_width, draw):
    lines, line = [], ""
    for ch in text:
        test = line + ch
        if draw.textlength(test, font=font) > max_width:
            lines.append(line)
            line = ch
        else:
            line = test
    if line:
        lines.append(line)
    return lines


def build_png_slides(slide_data, design, png_size=(1920, 1080)):
    """slide_dataとdesignからPNG画像リストを生成する。全サイズ対応（比例スケール）。"""
    W, H = png_size
    bg_c = hex_to_rgb(design["bg"])
    tc   = hex_to_rgb(design["title"])
    cc   = hex_to_rgb(design["content"])
    ac   = hex_to_rgb(design["accent"])

    telop_h = int(H * 0.22)
    area_h  = H - telop_h

    # 基準解像度 1920×1080 からのスケール係数
    sw = W / 1920
    sh = H / 1080
    s  = min(sw, sh)   # フォントは短辺に合わせる

    def fw(v): return max(1, int(v * sw))   # 横ピクセル
    def fh(v): return max(1, int(v * sh))   # 縦ピクセル
    def fs(v): return max(8, int(v * s))    # フォントサイズ

    f_huge    = get_font(fs(88))
    f_title   = get_font(fs(56))
    f_content = get_font(fs(36))
    f_small   = get_font(fs(28))

    lh_title   = fs(72)   # タイトル行間
    lh_content = fs(50)   # 本文行間

    def draw_centered(draw, text, font, lh, center_y, max_w, color):
        lines = wrap_text(text, font, max_w, draw)
        total = len(lines) * lh
        y = center_y - total // 2
        for line in lines:
            draw.text((W // 2, y), line, font=font, fill=color, anchor="mt")
            y += lh
        return y

    images = []

    # ── タイトルスライド ──
    img = Image.new("RGB", (W, H), bg_c)
    draw = ImageDraw.Draw(img)
    cy = area_h // 2
    draw.rectangle([0, cy - fh(6), W, cy + fh(6)], fill=ac)
    draw_centered(draw, slide_data.get("title", ""), f_huge,
                  fs(100), cy - fh(60), W - fw(120), tc)
    images.append(img)

    for slide_info in slide_data.get("slides", []):
        layout  = slide_info.get("layout", "standard")
        emoji   = slide_info.get("emoji", "")
        title   = slide_info.get("title", "")
        content = slide_info.get("content", [])

        img = Image.new("RGB", (W, H), bg_c)
        draw = ImageDraw.Draw(img)

        if layout == "impact":
            # ── インパクト：枠線＋中央大テキスト ──
            bw = max(4, fh(12))
            draw.rectangle([0, 0, W - 1, area_h - 1], outline=ac, width=bw)
            if emoji:
                draw.text((fw(60), fh(48)), emoji, font=f_title, fill=ac)
            has_sub = bool(content)
            cy = area_h // 2 - (fh(40) if has_sub else 0)
            bot = draw_centered(draw, title, f_huge, fs(100), cy, W - fw(160), tc)
            for item in (content[:2] if has_sub else []):
                draw.text((W // 2, bot + fh(12)), item, font=f_content, fill=cc, anchor="mt")
                bot += lh_content

        elif layout == "section":
            # ── セクション：左バー＋上下ライン＋中央タイトル ──
            draw.rectangle([0, 0, fw(18), H], fill=ac)
            draw.rectangle([0, 0, W, fh(10)], fill=ac)
            draw.rectangle([0, area_h - fh(10), W, area_h], fill=ac)
            emoji_off = 0
            if emoji:
                draw.text((W // 2, area_h // 2 - fh(100)),
                          emoji, font=f_huge, fill=ac, anchor="mm")
                emoji_off = fh(70)
            draw_centered(draw, title, f_title, lh_title,
                          area_h // 2 + emoji_off, W - fw(100), tc)
            if content:
                draw.text((W // 2, area_h // 2 + emoji_off + fh(90)),
                          content[0], font=f_small, fill=cc, anchor="mt")

        else:
            # ── standard：上バー＋タイトル折り返し＋箇条書き ──
            bar_h = fh(12)
            draw.rectangle([0, 0, W, bar_h], fill=ac)

            title_str = f"{emoji}  {title}" if emoji else title
            t_lines = wrap_text(title_str, f_title, W - fw(140), draw)
            y = fh(38)
            for line in t_lines:
                draw.text((fw(70), y), line, font=f_title, fill=tc)
                y += lh_title

            sep_y = y + fh(8)
            draw.rectangle([fw(70), sep_y, W - fw(70), sep_y + fh(4)], fill=ac)

            y = sep_y + fh(18)
            for item in content:
                for line in wrap_text(f"  •  {item}", f_content, W - fw(160), draw):
                    if y + lh_content > area_h:
                        break
                    draw.text((fw(80), y), line, font=f_content, fill=cc)
                    y += lh_content
                y += fh(8)

        images.append(img)

    return images


def build_png_zip(images, display_name):
    """PNG画像リストをZIPにまとめてバイナリを返す。"""
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, img in enumerate(images):
            img_buf = BytesIO()
            img.save(img_buf, format="PNG")
            zf.writestr(f"slide_{i+1:03d}.png", img_buf.getvalue())
    buf.seek(0)
    return buf.getvalue()


def group_slides_by_episode(slide_data):
    """スライドを episode フィールドでグループ分けして {ep_num: [slides]} を返す。"""
    groups = {}
    for s in slide_data.get("slides", []):
        ep = int(s.get("episode", 1))
        groups.setdefault(ep, []).append(s)
    return dict(sorted(groups.items()))


def build_slides_zip(slide_data, design, fmt, output_formats, display_name):
    """話数ごとにフォルダ分けし、PNG と PPTX を1つのZIPにまとめて返す。"""
    episodes = group_slides_by_episode(slide_data)
    video_title = slide_data.get("title", display_name)
    is_multi = len(episodes) > 1

    buf = BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for ep_num, slides in episodes.items():
            ep_label = f"第{ep_num}話" if is_multi else display_name
            ep_title = f"{video_title}　第{ep_num}話" if is_multi else video_title
            ep_data  = {"title": ep_title, "slides": slides}
            prefix   = f"{ep_label}/" if is_multi else ""

            if "PNG (.zip)" in output_formats:
                imgs = build_png_slides(ep_data, design, fmt["png"])
                for i, img in enumerate(imgs):
                    img_buf = BytesIO()
                    img.save(img_buf, format="PNG")
                    zf.writestr(f"{prefix}slide_{i+1:03d}.png", img_buf.getvalue())

            if "PPT (.pptx)" in output_formats:
                pptx_bytes = build_pptx(ep_data, design, fmt["pptx"])
                fname = f"{ep_label}.pptx" if is_multi else "slides.pptx"
                zf.writestr(f"{prefix}{fname}", pptx_bytes)

    buf.seek(0)
    return buf.getvalue()


def extract_script_only(script_text):
    """3カラムテーブルからセリフ列だけを抽出してプレーンテキストに変換する。
    テーブル形式でない台本はそのまま返す。"""
    lines = script_text.split('\n')
    if not any('|' in line for line in lines):
        return script_text
    result = []
    current_section = ""
    for line in lines:
        s = line.strip()
        if s.startswith('## ') or s.startswith('# '):
            current_section = s.lstrip('#').strip()
            result.append(f"\n【{current_section}】\n")
        elif re.match(r'^\|[-|: ]+\|$', s):
            continue  # 区切り行
        elif s.startswith('|'):
            cells = [c.strip() for c in s.split('|')[1:-1]]
            if len(cells) >= 2:
                script_cell = cells[1]
                if '演者のセリフ' in script_cell or 'セリフ' in script_cell:
                    continue  # ヘッダー行
                if script_cell:
                    result.append(script_cell)
        elif s and not s.startswith('|'):
            result.append(line)
    return '\n'.join(result)


def extract_slide_csv(script_text):
    """3カラムテーブルからセクション・映像テロップ・セリフをCSV（UTF-8 BOM）に変換する。"""
    import csv as _csv, io as _io
    lines = script_text.split('\n')
    rows = []
    current_section = ""
    for line in lines:
        s = line.strip()
        if s.startswith('## ') or s.startswith('# '):
            current_section = s.lstrip('#').strip()
        elif re.match(r'^\|[-|: ]+\|$', s):
            continue
        elif s.startswith('|'):
            cells = [c.strip() for c in s.split('|')[1:-1]]
            if len(cells) >= 2:
                if '映像' in cells[0] or 'テロップ' in cells[0]:
                    continue  # ヘッダー行
                video  = cells[0] if len(cells) > 0 else ""
                script = cells[1] if len(cells) > 1 else ""
                if video or script:
                    rows.append([current_section, video, script])
    buf = _io.StringIO()
    w = _csv.writer(buf)
    w.writerow(['セクション', '映像/テロップ指示', '演者のセリフ'])
    w.writerows(rows)
    return '﻿' + buf.getvalue()  # Excel対応 UTF-8 BOM


def save_script(script, product_name):
    OUTPUT_DIR.mkdir(exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_name = "".join(c if c.isalnum() or c in "-_" else "_" for c in product_name)
    output_path = OUTPUT_DIR / f"{timestamp}_{safe_name}.txt"
    output_path.write_text(script, encoding="utf-8")
    return output_path


# ── UI ──────────────────────────────────────────────────────────────────────

st.set_page_config(page_title="台本生成システム", page_icon="🎬", layout="wide")
st.title("🎬 プロダクトローンチ 台本生成システム")

api_key = os.environ.get("ANTHROPIC_API_KEY", "")
if not api_key:
    api_key = st.text_input("Anthropic APIキー", type="password", placeholder="sk-ant-...")
    if not api_key:
        st.warning("APIキーを入力してください")
        st.stop()

tavily_api_key = os.environ.get("TAVILY_API_KEY", "")

# セッション初回起動時にファイルからプリセットを読み込む
if "saved_presets" not in st.session_state:
    st.session_state["saved_presets"] = load_presets_from_file()

with st.sidebar:
    st.header("プリセット管理")
    st.caption("""
**使い方**
- **保存**：台本生成後に下部の「この入力内容を保存する」でプリセット名をつけて保存
- **呼び出し**：下のリストから選んで「呼び出す」をクリック
- **クリア**：「フォームをクリア」で入力欄を初期化（プリセットは削除されません）
- **エクスポート/インポート**：JSONファイルで他のPCへの持ち出しや共有が可能
""")

    if st.button("フォームをクリア", use_container_width=True, key="clear_form_btn"):
        clear_form_state()
        st.toast("入力内容をクリアしました")
        st.rerun()

    saved_presets = st.session_state.get("saved_presets", {})

    # ── 削除確認ダイアログ ──
    _pending_del = st.session_state.get("confirm_delete_name")
    if _pending_del and _pending_del in saved_presets:
        st.warning(f"「{_pending_del}」を本当に削除しますか？")
        _cy, _cn = st.columns(2)
        with _cy:
            if st.button("削除する", type="primary", use_container_width=True, key="confirm_del_yes"):
                if "recently_deleted" not in st.session_state:
                    st.session_state["recently_deleted"] = {}
                rd = st.session_state["recently_deleted"]
                rd[_pending_del] = saved_presets[_pending_del]
                if len(rd) > 5:
                    del rd[next(iter(rd))]
                del st.session_state["saved_presets"][_pending_del]
                save_presets_to_file(st.session_state["saved_presets"])
                st.session_state.pop("confirm_delete_name", None)
                st.toast(f"「{_pending_del}」を削除しました（削除済みから復元できます）")
                st.rerun()
        with _cn:
            if st.button("キャンセル", use_container_width=True, key="confirm_del_no"):
                st.session_state.pop("confirm_delete_name", None)
                st.rerun()
        st.divider()

    # ── 呼び出し ──
    if saved_presets:
        sel = st.selectbox("保存済みプリセット", ["── 選択 ──"] + list(saved_presets.keys()), key="preset_select")
        if st.button("呼び出す", use_container_width=True):
            if sel != "── 選択 ──":
                load_preset_to_session(saved_presets[sel])
                st.session_state.last_info = saved_presets[sel].copy()
                st.success(f"「{sel}」を読み込みました")
                st.rerun()

        # ── プリセット一覧・削除 ──
        with st.expander("プリセット一覧・削除"):
            st.caption("削除ボタンを押すと確認ダイアログが表示されます")
            for _pname in list(saved_presets.keys()):
                _col_n, _col_d = st.columns([3, 1])
                _col_n.text(_pname)
                if _col_d.button("削除", key=f"del_btn_{_pname}", use_container_width=True):
                    st.session_state["confirm_delete_name"] = _pname
                    st.rerun()
    else:
        st.caption("保存済みプリセットはありません")

    # ── 削除済みプリセット（復元） ──
    recently_deleted = st.session_state.get("recently_deleted", {})
    if recently_deleted:
        st.caption("削除済み（復元可能）")
        for del_name in list(recently_deleted.keys()):
            col_r, col_rb = st.columns([3, 1])
            col_r.text(del_name)
            if col_rb.button("復元", key=f"restore_{del_name}"):
                st.session_state["saved_presets"][del_name] = recently_deleted[del_name]
                save_presets_to_file(st.session_state["saved_presets"])
                del st.session_state["recently_deleted"][del_name]
                st.toast(f"「{del_name}」を復元しました")
                st.rerun()

    # JSONインポート
    st.caption("JSONファイルから読み込む")
    uploaded = st.file_uploader("インポート", type=["json"], label_visibility="collapsed")
    if uploaded:
        try:
            data = json.loads(uploaded.read())
            if "saved_presets" not in st.session_state:
                st.session_state["saved_presets"] = {}
            st.session_state["saved_presets"].update(data)
            save_presets_to_file(st.session_state["saved_presets"])
            st.success("インポートしました")
            st.rerun()
        except Exception:
            st.error("JSONの読み込みに失敗しました")

    # JSONエクスポート
    if saved_presets:
        export_json = json.dumps(saved_presets, ensure_ascii=False, indent=2)
        st.download_button(
            "プリセットをエクスポート (.json)",
            data=export_json.encode("utf-8"),
            file_name="presets.json",
            mime="application/json",
            use_container_width=True,
        )

    st.divider()
    st.header("サンプル台本")
    samples = load_samples()
    if samples:
        st.success(f"{len(samples)} 件読み込み済み")
        for s in samples:
            st.text(f"• {s['filename']}")
    else:
        st.warning("samplesフォルダにファイルがありません")
    if st.button("再読み込み", use_container_width=True):
        st.rerun()

if "system_blocks" not in st.session_state or st.session_state.get("samples_count") != len(samples):
    st.session_state.system_blocks = build_system_prompt(samples)
    st.session_state.samples_count = len(samples)

# ── フォーム ─────────────────────────────────────────────────────────────────

with st.form("product_form"):

    # 構成タイプ
    st.subheader("構成タイプ")
    structure_type = st.radio(
        "タイプを選択",
        ["従来型", "フロントエンド型"],
        captions=[
            "実績提示→商品説明→販売。インタビュー or 一人語り。スマートイット・アルゴテック型。",
            "無料ノウハウ・セミナー形式で価値提供→「時間・リスクが心配な方はこちら」と商品販売に自然につなぐ。",
        ],
        horizontal=True,
        key="f_structure_type",
    )

    st.divider()

    # フロントエンド型ノウハウパート
    st.subheader("フロントエンド型ノウハウパート")
    include_knowhow = st.checkbox("ノウハウパートを含める", value=False, key="f_include_knowhow")
    st.caption("オンにすると、指定テーマのノウハウ・価値提供コンテンツを台本に自動で組み込みます")
    col_kh1, col_kh2 = st.columns(2)
    with col_kh1:
        knowhow_theme = st.text_input(
            "ノウハウのテーマ / キーワード",
            placeholder="例：FXスキャルピング、副業で稼ぐ方法、仮想通貨の始め方",
            key="f_knowhow_theme",
        )
    with col_kh2:
        knowhow_notes = st.text_area(
            "ノウハウの補足メモ（任意）",
            placeholder="例：初心者向けに5分足を使った手法を説明したい、具体的なエントリーポイントを入れてほしい",
            height=80,
            key="f_knowhow_notes",
        )

    st.divider()

    # 商品・基本情報
    st.subheader("商品・基本情報")
    col1, col2 = st.columns(2)
    with col1:
        name = st.text_input("商品名", placeholder="例：スマートイット", key="f_name")
        category = st.selectbox("ジャンル", [
            "FX・為替投資", "株式投資・トレード", "仮想通貨・Web3",
            "副業・ビジネス", "美容・スキンケア", "健康・ダイエット",
            "教育・スキルアップ", "テック・SaaS", "食品・サプリ", "その他"
        ], key="f_category")
        seller_name = st.text_input("販売者名", placeholder="例：はたけ", key="f_seller_name")
        seller_first_person = st.selectbox("販売者の一人称", ["私", "僕", "おれ"], key="f_seller_first_person")
        seller_profile = st.text_area(
            "販売者のプロフィール",
            placeholder="例：元会社員で副業からFXを始め、3年で脱サラ。現在はFX系YouTuberとして活動中。フォロワー10万人。",
            height=80, key="f_seller_profile",
        )
        interviewer_name = st.text_input(
            "インタビュアー名（空欄で一人語り形式）",
            placeholder="例：ふじき　　※空欄→モノローグ形式",
            key="f_interviewer_name",
        )
        interviewer_first_person = st.selectbox("インタビュアーの一人称", ["私", "僕", "おれ"], key="f_interviewer_first_person")
        interviewer_profile = st.text_area(
            "インタビュアーのプロフィール（任意）",
            placeholder="例：元銀行員、現在は投資系メディアのライター。読者目線で質問するのが得意。",
            height=80, key="f_interviewer_profile",
        )
        seller_authority = st.text_input("販売者の権威・実績", placeholder="例：FX系YouTuber、動画1500本以上、5年以上活動", key="f_seller_authority")
        seller_story = st.text_area(
            "原体験・どん底ストーリー",
            placeholder="例：会社員時代にFXで300万円溶かし借金生活に。毎晩眠れない日々が続いた。そこで偶然ある手法に出会い、半年で完済。その経験を体系化したのがこのシステム。",
            height=110, key="f_seller_story",
        )
    with col2:
        catchcopy = st.text_input("キャッチコピー", placeholder="例：月5分で月10万円", key="f_catchcopy")
        target_audience = st.text_input("ターゲット層", placeholder="例：投資初心者、副業したい人", key="f_target_audience")
        result1 = st.text_input("実績数値①", placeholder="例：3ヶ月で31万円の利益", key="f_result1")
        result2 = st.text_input("実績数値②", placeholder="例：1年で125万円の利益", key="f_result2")
        monthly_return = st.text_input("月利 / 月収目安", placeholder="例：月利10%、月10万円", key="f_monthly_return")
        ease_of_start = st.text_input("始めやすさの根拠", placeholder="例：2万円から、スマホだけでOK", key="f_ease_of_start")
        mechanism = st.text_area(
            "独自メカニズム・勝てる理由",
            placeholder="例：2万通りの相場パターンをAIが学習。プロアナリストが検証した上位3%のシグナルのみを抽出して自動売買。完全放置でも機能する理由はここにある。",
            height=110, key="f_mechanism",
        )

    st.markdown("**商品の強み**（入力した分だけ使用）")
    s_cols = st.columns(4)
    strength_placeholders = [
        "例：証拠金2万円から始められる",
        "例：勝率78%で安心感がある",
        "例：ドローダウン平均10%以下",
        "例：完全自動、ON/OFFも不要",
    ]
    strengths = []
    for i, col in enumerate(s_cols):
        s = col.text_input(f"強み{i+1}", key=f"str_{i}", label_visibility="collapsed", placeholder=strength_placeholders[i])
        strengths.append(s)

    st.divider()

    # 利用者の声
    st.subheader("利用者の声（喜びの声）")
    st.caption("名前・属性と声のセットで入力してください")
    voice1 = st.text_input("声①", placeholder="例：30代会社員Aさん：導入1ヶ月で8万円の利益が出ました！初心者でも全く問題なかったです。", key="f_voice1")
    voice2 = st.text_input("声②", placeholder="例：40代主婦Bさん：スマホだけで設定できて、今は毎月安定して収入が入っています。", key="f_voice2")
    voice3 = st.text_input("声③", placeholder="例：20代フリーランスCさん：他のシステムで失敗したけどこれは違いました。", key="f_voice3")

    st.divider()

    # 社会背景・痛み訴求
    st.subheader("社会背景・痛み訴求")
    col3, col4 = st.columns(2)
    with col3:
        pain_points = st.text_area("視聴者のペイン", placeholder="例：残業しても給料が上がらない、将来が不安、副業する時間がない", height=100, key="f_pain_points")
    with col4:
        why_now = st.text_area("なぜ今必要か（why now）", placeholder="例：物価は上がるのに賃金は上がらない。自分で資産を作るしかない。", height=100, key="f_why_now")
    col5, col6 = st.columns(2)
    with col5:
        vs_competition = st.text_area(
            "競合との違い・失敗体験との対比",
            placeholder="例：他の自動売買は設定が複雑で初心者が挫折しやすい。このシステムは設定を代行してくれるので、届いた当日から稼働できる。",
            height=100, key="f_vs_competition",
        )
    with col6:
        before_after = st.text_area(
            "購入後のビフォーアフター",
            placeholder="例：Before：毎朝5時起きでチャートを見ながら出勤。After：朝ゆっくり起きてスマホで確認するだけ。家族との時間が増えた。",
            height=100, key="f_before_after",
        )

    st.divider()

    # 第三者・信頼性
    st.subheader("第三者・信頼性")
    col5, col6, col7 = st.columns(3)
    with col5:
        third_party_type = st.selectbox("第三者の種類", ["なし", "開発者・専門家", "ユーザー・実践者", "有識者・研究者", "著名人・インフルエンサー"], key="f_third_party_type")
    with col6:
        third_party_name = st.text_input("名前・肩書き", placeholder="例：桑田（システム開発者）", key="f_third_party_name")
    with col7:
        third_party_points = st.text_input("第三者の裏付けポイント", placeholder="例：2万通りのロジックを検証、14歳からシステム開発", key="f_third_party_points")

    st.divider()

    # 販売フロー
    st.subheader("販売フロー")
    sales_flow_type = st.radio(
        "販売フロータイプを選択",
        ["直接販売型", "面談・相談誘導型"],
        captions=[
            "ローンチ動画から直接購入ページへ誘導して販売する",
            "個別Zoom面談・LINE相談などに誘導し、個別接触後に販売する",
        ],
        horizontal=True,
        key="f_sales_flow_type",
    )
    col_sf1, col_sf2 = st.columns(2)
    with col_sf1:
        sales_start_day = st.selectbox(
            "販売スタート日（直接販売型）",
            [
                "当日（即時販売）",
                "翌日（1日後）",
                "翌々日（2日後）",
                "3日後",
                "4日後",
                "5日後",
                "7日後（1週間後）",
                "その他（追加メモに記載）",
            ],
            index=1,
            key="f_sales_start_day",
            help="直接販売型を選択した場合に使用されます",
        )
    with col_sf2:
        consultation_method = st.text_input(
            "誘導先の詳細（面談・相談誘導型）",
            placeholder="例：個別Zoom面談、LINE@で個別相談",
            key="f_consultation_method",
            help="面談・相談誘導型を選択した場合に使用されます",
        )

    st.divider()

    # 価格・オファー
    st.subheader("価格・オファー")
    col8, col9, col10 = st.columns(3)
    with col8:
        regular_price = st.text_input("定価", placeholder="例：158,000円", key="f_regular_price")
        special_price = st.text_input("特別価格", placeholder="例：98,000円", key="f_special_price")
    with col9:
        limited_time = st.text_input("期間限定条件", placeholder="例：3日間限定", key="f_limited_time")
        limited_seats = st.text_input("先着・定員制限（任意）", placeholder="例：先着50名様、定員30名", key="f_limited_seats")
        installment = st.selectbox("分割対応", ["なし", "あり"], key="f_installment")
    with col10:
        bonuses = st.text_area("特典内容（カンマ区切り）", placeholder="例：導入マニュアル、トレーダー手法、キャッシュを増やす方法", height=100, key="f_bonuses")

    st.divider()

    # コメント促進パート
    st.subheader("コメント促進パート（動画末尾）")
    st.caption("話ごとに含める/含めないを選択できます。質問は空欄で自動生成します。")
    comment_includes = []
    comment_prompts = []
    for i in range(5):
        col_a, col_b = st.columns([1, 4])
        with col_a:
            include = st.checkbox(f"第{i+1}話", key=f"comment_include_{i}", value=True)
        with col_b:
            cp = st.text_input(
                f"第{i+1}話の質問",
                key=f"comment_{i}",
                label_visibility="collapsed",
                placeholder="質問を入力（空欄で自動生成）"
            )
        comment_includes.append(include)
        comment_prompts.append(cp)

    st.divider()

    # 構成オプション
    st.subheader("構成オプション")
    col11, col12 = st.columns(2)
    with col11:
        episode_structure = st.selectbox("話数構成", [
            "1話完結", "2話構成（前編・後編）", "3話構成", "4話構成", "5話構成"
        ], key="f_episode_structure")
        closing_strength = st.selectbox("クロージングの強度", [
            "真摯・控えめ（押し付けない）", "標準（バランス型）", "強め（urgency高め）", "最強（限定・希少性全開）"
        ], key="f_closing_strength")
        video_duration = st.selectbox("1話あたりの動画の長さ", [
            "3分（約900文字）",
            "5分（約1,500文字）",
            "7分（約2,100文字）",
            "10分（約3,000文字）",
            "15分（約4,500文字）",
            "20分（約6,000文字）",
            "30分（約9,000文字）",
            "45分（約13,500文字）",
            "60分（約18,000文字）",
            "90分（約27,000文字）",
            "120分（約36,000文字）",
        ], index=2, key="f_video_duration")
    with col12:
        notes = st.text_area("追加メモ（任意）", placeholder="例：競合との比較を入れたい、このワードは避けたいなど", height=100, key="f_notes")

    st.divider()

    # 話数ごとの内容指定
    st.subheader("話数ごとの内容指定")
    use_episode_themes = st.checkbox("各話の内容・テーマを指定する", value=False, key="f_use_episode_themes")
    st.caption("オンにすると、各話で話す内容を個別に指定できます。空欄の話はAIが自動で構成します。")
    episode_themes = []
    for i in range(5):
        theme = st.text_area(
            f"第{i+1}話の内容・テーマ",
            placeholder=f"例：第{i+1}話では〇〇について話す。ポイントは△△と□□。",
            height=80,
            key=f"episode_theme_{i}",
        )
        episode_themes.append(theme)

    st.divider()

    st.markdown("**出力形式を選択**")
    output_format_normal = st.radio(
        "出力形式",
        ["台本のみ（セリフのみ・API費用削減）", "完全版（映像指示+セリフ+演技指示）"],
        key="normal_output_format",
        horizontal=True,
        label_visibility="collapsed",
    )

    use_trend_search = st.checkbox(
        "最新トレンドをWeb検索して台本に反映する",
        value=False,
        disabled=not tavily_api_key,
        help="TavilyのAPIキーが設定されている場合に利用できます",
    )

    # ── プリセット保存（フォーム内）──
    st.divider()
    st.caption("📌 入力内容をプリセット保存する（保存ボタンを押すと台本は生成されません）")
    col_pn2, col_pb2 = st.columns([3, 1])
    with col_pn2:
        preset_save_name_form = st.text_input(
            "プリセット名",
            placeholder="例：スマートイット用、FX商品A用",
            key="preset_save_name_form",
        )
    with col_pb2:
        st.write("")
        save_preset_submitted = st.form_submit_button("保存のみ", use_container_width=True)

    submitted = st.form_submit_button("台本を生成する（通常モード）", type="primary", use_container_width=True)

# ── プリセット保存処理（フォーム送信後）──────────────────────────────────────

if save_preset_submitted:
    preset_name_val = st.session_state.get("preset_save_name_form", "").strip()
    if preset_name_val:
        if "saved_presets" not in st.session_state:
            st.session_state["saved_presets"] = {}
        _save_info = build_info_from_session()
        existing = st.session_state["saved_presets"]
        save_name = preset_name_val
        if save_name in existing:
            n = 2
            while f"{preset_name_val}（{n}）" in existing:
                n += 1
            save_name = f"{preset_name_val}（{n}）"
        st.session_state["saved_presets"][save_name] = _save_info
        st.session_state.last_info = _save_info
        save_presets_to_file(st.session_state["saved_presets"])
        if save_name != preset_name_val:
            st.success(f"「{preset_name_val}」が既に存在するため「{save_name}」として保存しました。")
        else:
            st.success(f"「{save_name}」を保存しました。左のサイドバーから呼び出せます。")
    else:
        st.warning("プリセット名を入力してください")

# ── 生成処理 ──────────────────────────────────────────────────────────────────

DURATION_MAX_TOKENS = {
    "3分": 2048,
    "5分": 3000,
    "7分": 4096,
    "10分": 6000,
    "15分": 8192,
    "20分": 12000,
    "30分": 16000,
    "45分": 20000,
    "60分": 24000,
    "90分": 28000,
    "120分": 32000,
}


def run_generation(client, system_blocks, messages, display_name, max_tokens=4096):
    """Claudeにリクエストを送りストリーミング表示する。台本とstatsを返す。"""
    script = ""
    placeholder = st.empty()
    with client.messages.stream(
        model=MODEL,
        max_tokens=max_tokens,
        system=system_blocks,
        messages=messages,
    ) as stream:
        for text in stream.text_stream:
            script += text
            placeholder.markdown(script)
        final = stream.get_final_message()
    usage = final.usage
    stats = {
        "input_tokens": usage.input_tokens,
        "output_tokens": usage.output_tokens,
        "cache_creation_tokens": getattr(usage, "cache_creation_input_tokens", 0),
        "cache_read_tokens": getattr(usage, "cache_read_input_tokens", 0),
    }
    return script, stats


def show_download(script, display_name, stats):
    """統計情報とダウンロードボタンを表示する。"""
    output_path = save_script(script, display_name)
    st.divider()
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("入力トークン", f"{stats['input_tokens']:,}")
    c2.metric("出力トークン", f"{stats['output_tokens']:,}")
    if stats["cache_creation_tokens"]:
        c3.metric("キャッシュ書込み", f"{stats['cache_creation_tokens']:,}")
    if stats["cache_read_tokens"]:
        c4.metric("キャッシュ読込み", f"{stats['cache_read_tokens']:,}")
    st.success(f"保存済み: {output_path}")
    st.download_button(
        "台本をダウンロード (.txt)",
        data=script.encode("utf-8"),
        file_name=f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{display_name}.txt",
        mime="text/plain",
        use_container_width=True,
        key=f"dl_{datetime.now().strftime('%H%M%S%f')}",
    )


def render_download_buttons(content, dl_name, ts, key_suffix):
    """3種類のダウンロードボタンを横並びで表示するヘルパー。"""
    dc1, dc2, dc3 = st.columns(3)
    with dc1:
        st.download_button(
            "完全版 (.txt)",
            data=content.encode("utf-8"),
            file_name=f"{ts}_{dl_name}.txt",
            mime="text/plain",
            use_container_width=True,
            key=f"dl_md{key_suffix}",
            help="映像指示・セリフ・演技指示をすべて含む完全版テキスト",
        )
    with dc2:
        st.download_button(
            "台本のみ (.txt)",
            data=extract_script_only(content).encode("utf-8"),
            file_name=f"{ts}_{dl_name}_セリフのみ.txt",
            mime="text/plain",
            use_container_width=True,
            key=f"dl_txt{key_suffix}",
            help="セリフ列だけを抽出したプレーンテキスト。読み合わせや練習用",
        )
    with dc3:
        st.download_button(
            "映像+台本 (.csv)　スライド用",
            data=extract_slide_csv(content).encode("utf-8"),
            file_name=f"{ts}_{dl_name}_slide.csv",
            mime="text/csv",
            use_container_width=True,
            key=f"dl_csv{key_suffix}",
            help="セクション・映像テロップ指示・セリフの2列CSV。スライド壁打ちに最適",
        )


if submitted:
    info = {
        "structure_type": structure_type,
        "include_knowhow": include_knowhow,
        "knowhow_theme": knowhow_theme,
        "knowhow_notes": knowhow_notes,
        "name": name, "category": category,
        "seller_name": seller_name, "seller_first_person": seller_first_person, "seller_profile": seller_profile,
        "interviewer_name": interviewer_name, "interviewer_first_person": interviewer_first_person, "interviewer_profile": interviewer_profile,
        "seller_authority": seller_authority,
        "catchcopy": catchcopy, "target_audience": target_audience,
        "result1": result1, "result2": result2,
        "monthly_return": monthly_return, "ease_of_start": ease_of_start,
        "strengths": strengths,
        "voices": [voice1, voice2, voice3],
        "pain_points": pain_points, "why_now": why_now,
        "vs_competition": vs_competition, "before_after": before_after,
        "seller_story": seller_story, "mechanism": mechanism,
        "third_party_type": third_party_type, "third_party_name": third_party_name,
        "third_party_points": third_party_points,
        "regular_price": regular_price, "special_price": special_price,
        "limited_time": limited_time, "limited_seats": limited_seats,
        "installment": installment, "bonuses": bonuses,
        "comment_includes": comment_includes,
        "comment_prompts": comment_prompts,
        "episode_structure": episode_structure, "closing_strength": closing_strength,
        "video_duration": video_duration,
        "use_episode_themes": use_episode_themes,
        "episode_themes": episode_themes,
        "notes": notes,
        "sales_flow_type": sales_flow_type,
        "sales_start_day": sales_start_day,
        "consultation_method": consultation_method,
    }
    _fmt_raw_n = st.session_state.get("normal_output_format", "完全版（映像指示+セリフ+演技指示）")
    _output_fmt_n = "台本のみ" if "台本のみ" in _fmt_raw_n else "完全版"
    user_prompt = build_user_prompt(info, _output_fmt_n)
    display_name = name if name else "台本"

    # トレンド検索
    if use_trend_search and tavily_api_key:
        with st.spinner("最新トレンドを検索中..."):
            trends = search_trends(tavily_api_key, info.get("category", ""), info.get("name", ""))
        if trends:
            user_prompt += f"\n\n## 最新トレンド・時事情報（Web検索結果）\n{trends}\n\n上記のトレンド情報も台本に自然に盛り込んでください。"
            st.info("最新トレンドを取得しました。台本に反映します。")

    st.divider()
    st.subheader("生成結果")

    try:
        client = anthropic.Anthropic(api_key=api_key)
        # 動画の長さからmax_tokensを決定
        duration_key = info.get("video_duration", "7分（約2,100文字）").split("（")[0]
        episode_num = int(info.get("episode_structure", "1話完結")[0]) if info.get("episode_structure", "1")[0].isdigit() else 1
        base_tokens = DURATION_MAX_TOKENS.get(duration_key, 4096)
        max_tokens = min(base_tokens * episode_num, 32000)

        messages = [{"role": "user", "content": user_prompt}]
        script, stats = run_generation(client, st.session_state.system_blocks, messages, display_name, max_tokens)

        # セッションに保存（再編集用）
        st.session_state.current_script = script
        st.session_state.current_messages = messages
        st.session_state.display_name = display_name
        st.session_state.last_info = info
        st.session_state.last_stats = stats

    except anthropic.APIError as e:
        st.error(f"APIエラー: {e}")


# ── 高精度モード（2ステップ生成） ─────────────────────────────────────────────

st.divider()
with st.expander("高精度モード（2ステップ生成）", expanded=("outline_draft" in st.session_state or st.session_state.get("block_gen_active", False))):
    st.caption("① 構成案を生成 → 確認・編集 → ② 台本を生成。意図に沿った高品質な台本ができます。")

    # ── ブロック分割生成中の処理 ──
    if st.session_state.get("block_gen_active", False):
        _bq = st.session_state.get("block_gen_queue", [])
        _bc = st.session_state.get("block_gen_completed", [])
        _bt = st.session_state.get("block_gen_total", 1)
        _bdn = st.session_state.get("block_gen_display_name", "台本")
        _bi = _bt - len(_bq)  # 完了済みブロック数

        _prog_col, _stop_col = st.columns([4, 1])
        with _prog_col:
            st.progress(_bi / _bt, text=f"{_bi}/{_bt} ブロック完了")
        with _stop_col:
            _stop_clicked = st.button("⏹ 中止", key="stop_block_gen")

        if _stop_clicked:
            st.session_state.block_gen_active = False
            if _bc:
                _partial = '\n\n---\n\n'.join(_bc)
                st.session_state.current_script = _partial
                st.session_state.display_name = _bdn
                st.session_state.last_info = st.session_state.get("block_gen_info", {})
                st.session_state.last_stats = {}
                save_script(_partial, _bdn)
            st.warning("⏹ 生成を中止しました。生成済みのブロックを保存しました。")
            st.rerun()

        if _bq:
            _block = _bq[0]
            st.markdown(f"**生成中：{_block['title']}（{_bi+1}/{_bt}）**")
            _prev = '\n\n'.join(_bc)
            _blk_prompt = build_block_script_prompt(
                st.session_state.get("block_gen_outline", ""),
                _block, _bi + 1, _bt, _prev,
                st.session_state.get("block_gen_info", {}),
                st.session_state.get("block_gen_output_format", "完全版"),
            )
            _blk_text = ""
            _blk_ph = st.empty()
            try:
                _client_blk = anthropic.Anthropic(api_key=api_key)
                with _client_blk.messages.stream(
                    model=MODEL,
                    max_tokens=6000,
                    system=st.session_state.system_blocks,
                    messages=[{"role": "user", "content": _blk_prompt}],
                ) as _stream:
                    for _t in _stream.text_stream:
                        _blk_text += _t
                        _blk_ph.markdown(_blk_text)
                # 話（エピソード）が変わったときにヘッダーを付与
                _cur_ep = _block.get('episode_header', '')
                _prev_ep = st.session_state.get("block_gen_current_episode", "")
                if _cur_ep and _cur_ep != _prev_ep:
                    _ep_prefix = f"# {_cur_ep}\n\n"
                    st.session_state.block_gen_current_episode = _cur_ep
                else:
                    _ep_prefix = ""
                # _block['title'] はすでに "## [N] ..." 形式なので ## を付け足さない
                _new_bc = _bc + [f"{_ep_prefix}{_block['title']}\n\n{_blk_text}"]
                st.session_state.block_gen_completed = _new_bc
                st.session_state.block_gen_queue = _bq[1:]
                if len(_bq) == 1:  # 最後のブロック
                    st.session_state.block_gen_active = False
                    _final_script = '\n\n---\n\n'.join(_new_bc)
                    st.session_state.current_script = _final_script
                    st.session_state.display_name = _bdn
                    st.session_state.last_info = st.session_state.get("block_gen_info", {})
                    st.session_state.last_stats = {}
                    save_script(_final_script, _bdn)
                st.rerun()
            except anthropic.APIError as _e:
                st.error(f"APIエラー: {_e}")
                st.session_state.block_gen_active = False

    st.markdown("**出力形式を選択**")
    hq_output_format = st.radio(
        "出力形式",
        ["台本のみ（セリフのみ・API費用削減）", "完全版（映像指示+セリフ+演技指示）"],
        key="hq_output_format",
        horizontal=True,
        label_visibility="collapsed",
    )

    hq_trend = st.checkbox(
        "最新トレンドをWeb検索して構成案に反映する",
        value=False,
        key="hq_use_trend",
        disabled=not tavily_api_key,
        help="TavilyのAPIキーが設定されている場合に利用できます",
    )

    col_hq1, col_hq2 = st.columns(2)
    with col_hq1:
        gen_outline_btn = st.button("① 構成案を生成する", use_container_width=True, key="btn_gen_outline")
    with col_hq2:
        if "outline_draft" in st.session_state:
            if st.button("構成案をクリアする", use_container_width=True, key="btn_clear_outline"):
                st.session_state.pop("outline_draft", None)
                st.session_state.pop("outline_info", None)
                st.rerun()

    if gen_outline_btn:
        if not api_key:
            st.error("APIキーを設定してください（サイドバー）")
        else:
            outline_info = build_info_from_session()
            outline_prompt = build_outline_prompt(outline_info)
            if hq_trend and tavily_api_key:
                with st.spinner("最新トレンドを検索中..."):
                    trends = search_trends(tavily_api_key, outline_info.get("category", ""), outline_info.get("name", ""))
                if trends:
                    outline_prompt += f"\n\n## 最新トレンド・時事情報（Web検索結果）\n{trends}\n\n上記のトレンド情報も構成案に自然に盛り込んでください。"
                    st.info("最新トレンドを取得しました。構成案に反映します。")
            episode_num = int(outline_info.get("episode_structure", "1話完結")[0]) if outline_info.get("episode_structure", "1")[0].isdigit() else 1
            outline_tokens = min(5000 * episode_num, 24000)
            st.markdown("**構成案を生成中...**")
            placeholder_o = st.empty()
            outline_text = ""
            try:
                client_o = anthropic.Anthropic(api_key=api_key)
                with client_o.messages.stream(
                    model=MODEL,
                    max_tokens=outline_tokens,
                    system=st.session_state.system_blocks,
                    messages=[{"role": "user", "content": outline_prompt}],
                ) as stream:
                    for text in stream.text_stream:
                        outline_text += text
                        placeholder_o.markdown(outline_text)
                st.session_state["outline_draft"] = outline_text
                st.session_state["outline_info"] = outline_info
                st.rerun()
            except anthropic.APIError as e:
                st.error(f"APIエラー: {e}")

    if "outline_draft" in st.session_state:
        st.markdown("**構成案（自由に編集できます）**")
        st.caption("内容・順序を確認・修正したら ② を押して台本を生成してください")
        outline_edit = st.text_area(
            "構成案",
            value=st.session_state["outline_draft"],
            height=520,
            key="outline_edit",
            label_visibility="collapsed",
        )
        gen_from_outline_btn = st.button(
            "② この構成で台本を生成する",
            type="primary",
            use_container_width=True,
            key="btn_gen_from_outline",
        )
        if gen_from_outline_btn:
            if not api_key:
                st.error("APIキーを設定してください（サイドバー）")
            else:
                outline_info2 = st.session_state.get("outline_info", build_info_from_session())
                display_name2 = (outline_info2.get("name", "台本") or "台本") + "_高精度"
                duration_key2 = outline_info2.get("video_duration", "7分（約2,100文字）").split("（")[0]
                episode_num2 = int(outline_info2.get("episode_structure", "1話完結")[0]) if outline_info2.get("episode_structure", "1")[0].isdigit() else 1

                # 15分以上かつブロックが2つ以上あればブロック分割生成
                use_block_mode = duration_key2 in ["15分", "20分", "30分", "45分", "60分", "90分", "120分"]
                blocks = parse_outline_blocks(outline_edit) if use_block_mode else []
                use_block_mode = use_block_mode and len(blocks) >= 2

                _fmt_raw = st.session_state.get("hq_output_format", "完全版（映像指示+セリフ+演技指示）")
                _output_fmt = "台本のみ" if "台本のみ" in _fmt_raw else "完全版"

                if use_block_mode:
                    # ブロック分割生成：キューを初期化してrerunで1ブロックずつ処理
                    st.session_state.block_gen_active = True
                    st.session_state.block_gen_queue = blocks
                    st.session_state.block_gen_completed = []
                    st.session_state.block_gen_total = len(blocks)
                    st.session_state.block_gen_outline = outline_edit
                    st.session_state.block_gen_info = outline_info2
                    st.session_state.block_gen_display_name = display_name2
                    st.session_state.block_gen_current_episode = ""
                    st.session_state.block_gen_output_format = _output_fmt
                    st.rerun()
                else:
                    try:
                        client2 = anthropic.Anthropic(api_key=api_key)
                        # 短尺：一括生成
                        script_prompt2 = build_script_from_outline_prompt(outline_edit, outline_info2)
                        base_tokens2 = DURATION_MAX_TOKENS.get(duration_key2, 4096)
                        max_tokens2 = min(base_tokens2 * episode_num2, 32000)
                        st.markdown("**台本を生成中...**")
                        placeholder2 = st.empty()
                        script2 = ""
                        with client2.messages.stream(
                            model=MODEL,
                            max_tokens=max_tokens2,
                            system=st.session_state.system_blocks,
                            messages=[{"role": "user", "content": script_prompt2}],
                        ) as stream:
                            for text in stream.text_stream:
                                script2 += text
                                placeholder2.markdown(script2)
                        st.session_state.current_script = script2
                        st.session_state.display_name = display_name2
                        st.session_state.last_info = outline_info2
                        st.session_state.last_stats = {}
                        save_script(script2, display_name2)
                    except anthropic.APIError as e:
                        st.error(f"APIエラー: {e}")


# ── 再編集パネル ──────────────────────────────────────────────────────────────

if "current_script" in st.session_state:

    # ── 台本ダウンロード（常時表示）──
    st.divider()
    st.subheader("生成済み台本")
    stats_s = st.session_state.get("last_stats", {})
    if stats_s:
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("入力トークン", f"{stats_s.get('input_tokens', 0):,}")
        c2.metric("出力トークン", f"{stats_s.get('output_tokens', 0):,}")
        if stats_s.get("cache_creation_tokens"):
            c3.metric("キャッシュ書込み", f"{stats_s['cache_creation_tokens']:,}")
        if stats_s.get("cache_read_tokens"):
            c4.metric("キャッシュ読込み", f"{stats_s['cache_read_tokens']:,}")
    _dl_name = st.session_state.get("display_name", "台本")
    _ts      = datetime.now().strftime('%Y%m%d_%H%M%S')
    _script  = st.session_state.current_script

    _episodes = split_script_by_episode(_script)

    if len(_episodes) <= 1:
        # 1話完結：タブなし
        render_download_buttons(_script, _dl_name, _ts, "_all")
    else:
        # 複数話：タブで話ごとに切り替え
        _tab_labels = ["全話"] + [f"第{i+1}話" for i in range(len(_episodes))]
        _dl_tabs = st.tabs(_tab_labels)
        with _dl_tabs[0]:
            render_download_buttons(_script, _dl_name, _ts, "_all")
        for _ei, (_tab, _ep) in enumerate(zip(_dl_tabs[1:], _episodes)):
            with _tab:
                st.caption(_ep['title'])
                render_download_buttons(
                    _ep['content'],
                    f"{_dl_name}_第{_ei+1}話",
                    _ts,
                    f"_ep{_ei+1}",
                )

    st.divider()
    st.subheader("再編集")
    st.caption("生成した台本に修正指示を出して再生成できます")

    edit_instruction = st.text_area(
        "修正指示",
        placeholder="例：クロージングをもっと強くして\n例：オープニングの問いかけを変えて\n例：〇〇のセクションを削除して",
        height=100,
        key="edit_instruction",
    )

    if st.button("再編集する", type="primary", use_container_width=True):
        if not edit_instruction.strip():
            st.warning("修正指示を入力してください")
        else:
            try:
                client = anthropic.Anthropic(api_key=api_key)
                display_name = st.session_state.display_name

                messages = st.session_state.current_messages + [
                    {"role": "assistant", "content": st.session_state.current_script},
                    {"role": "user", "content": f"以下の修正指示に従って台本を修正してください：\n\n{edit_instruction}"},
                ]

                st.subheader("修正結果")
                script, stats = run_generation(client, st.session_state.system_blocks, messages, display_name)

                st.session_state.current_script = script
                st.session_state.current_messages = messages
                st.session_state.last_stats = stats

            except anthropic.APIError as e:
                st.error(f"APIエラー: {e}")

    # ── スライド生成 ──────────────────────────────────────────────────────────
    st.divider()
    st.subheader("スライド作成")

    # スライドフォーマット・デザイン設定
    col_f, col_d1 = st.columns(2)
    with col_f:
        slide_format = st.selectbox("スライドフォーマット", list(SLIDE_FORMATS.keys()))
    with col_d1:
        design_preset = st.selectbox("デザインプリセット", list(DESIGN_PRESETS.keys()) + ["カスタム"])
    if design_preset == "カスタム":
        dc1, dc2, dc3, dc4 = st.columns(4)
        bg_c    = dc1.color_picker("背景色",      "#1a1a1a")
        title_c = dc2.color_picker("タイトル色",  "#ffffff")
        cont_c  = dc3.color_picker("テキスト色",  "#e0e0e0")
        acc_c   = dc4.color_picker("アクセント色", "#4a9eff")
        design = {"bg": bg_c, "title": title_c, "content": cont_c, "accent": acc_c}
    else:
        design = DESIGN_PRESETS[design_preset]

    output_formats = st.multiselect(
        "出力形式", ["PPT (.pptx)", "PNG (.zip)"], default=["PPT (.pptx)"],
    )

    if st.button("スライドを作成する", use_container_width=True):
        try:
            client = anthropic.Anthropic(api_key=api_key)
            chunks = split_script_into_chunks(st.session_state.current_script)
            total  = len(chunks)

            # 分割内容を事前に表示
            if total > 1:
                labels = " / ".join(lbl for _, lbl, _ in chunks)
                st.info(f"台本を **{total} 分割** して順番に生成します：{labels}")

            progress = st.progress(0)
            status   = st.empty()

            all_slides  = []
            video_title = ""

            for i, (ep_num, label, chunk_text) in enumerate(chunks):
                status.text(f"生成中：{label}（{i + 1} / {total}）")
                chunk_data = generate_slide_data(client, chunk_text)

                if not video_title:
                    video_title = chunk_data.get("title", "")

                for slide in chunk_data.get("slides", []):
                    slide["episode"] = ep_num
                    all_slides.append(slide)

                progress.progress((i + 1) / total)

            status.text("✅ 全チャンク生成完了！")
            slide_data = {"title": video_title, "slides": all_slides}

            st.session_state.slide_data           = slide_data
            st.session_state.slide_design         = design
            st.session_state.slide_format         = slide_format
            st.session_state.slide_output_formats = output_formats
            st.session_state.slide_preview_ready  = True
            st.session_state.slide_approved       = False
            for k in ("slide_pptx_bytes", "slide_zip_bytes", "slide_output_ts"):
                st.session_state.pop(k, None)
            st.rerun()
        except Exception as e:
            st.error(f"スライド生成エラー: {e}")

    # ── プレビュー & 承認フロー ──
    if st.session_state.get("slide_preview_ready"):
        _sd      = st.session_state.slide_data
        _design  = st.session_state.slide_design
        _fmt     = SLIDE_FORMATS[st.session_state.slide_format]
        _ofmts   = st.session_state.get("slide_output_formats", ["PPT (.pptx)"])
        _count   = len(_sd.get("slides", []))

        if not st.session_state.get("slide_approved"):
            # ── 全スライドプレビュー（サムネイル一覧）──
            st.info(f"全 **{_count} 枚**を確認してから出力してください。文字ズレや内容がおかしいスライドがないか確認してください。")

            orig_w, orig_h = _fmt["png"]
            prev_size = (orig_w // 4, orig_h // 4)

            with st.spinner(f"全スライドのプレビューを生成中（{_count} 枚）..."):
                prev_imgs = build_png_slides(_sd, _design, prev_size)

            n_cols = 4
            labels = ["タイトル"] + [f"スライド {i+1}" for i in range(len(prev_imgs) - 1)]
            for row_start in range(0, len(prev_imgs), n_cols):
                row = prev_imgs[row_start:row_start + n_cols]
                row_labels = labels[row_start:row_start + n_cols]
                cols = st.columns(n_cols)
                for col, img, lbl in zip(cols, row, row_labels):
                    col.image(img, caption=lbl, use_container_width=True)

            st.divider()
            col_ok, col_ng = st.columns(2)
            with col_ok:
                if st.button("OK、全部出力する", type="primary", use_container_width=True, key="slide_ok_btn"):
                    st.session_state.slide_approved = True
                    st.rerun()
            with col_ng:
                if st.button("やり直す（再生成）", use_container_width=True, key="slide_retry_btn"):
                    for k in ("slide_preview_ready", "slide_approved", "slide_data",
                              "slide_pptx_bytes", "slide_zip_bytes"):
                        st.session_state.pop(k, None)
                    st.rerun()

        else:
            # ── 全スライド出力（初回のみ生成してキャッシュ）──
            display_name = st.session_state.display_name
            ts = st.session_state.setdefault("slide_output_ts", datetime.now().strftime("%Y%m%d_%H%M%S"))

            if "slide_zip_bytes" not in st.session_state:
                episodes = group_slides_by_episode(_sd)
                ep_count = len(episodes)
                label = f"{_count} 枚 / {ep_count} 話" if ep_count > 1 else f"{_count} 枚"
                with st.spinner(f"スライドを出力中（{label}）..."):
                    st.session_state.slide_zip_bytes = build_slides_zip(
                        _sd, _design, _fmt, _ofmts, display_name
                    )

            episodes = group_slides_by_episode(_sd)
            is_multi = len(episodes) > 1
            if is_multi:
                ep_labels = "・".join(f"第{ep}話" for ep in episodes)
                st.info(f"フォルダ構成: {ep_labels}")

            st.download_button(
                "スライドをダウンロード (.zip)",
                data=st.session_state.slide_zip_bytes,
                file_name=f"{ts}_{display_name}_slides.zip",
                mime="application/zip",
                use_container_width=True,
                key="dl_slides_main",
            )
            st.success(f"全 {_count} 枚の出力が完了しました")

            if st.button("スライドを作り直す", use_container_width=True, key="slide_redo_btn"):
                for k in ("slide_preview_ready", "slide_approved",
                          "slide_zip_bytes", "slide_output_ts"):
                    st.session_state.pop(k, None)
                st.rerun()

    # ── スライド修正 ──────────────────────────────────────────────────────────
    if "slide_data" in st.session_state:
        st.divider()
        st.subheader("スライド修正")
        st.caption("生成したスライドに修正指示を出して再生成できます")

        slide_edit = st.text_area(
            "修正指示",
            placeholder="例：スライド3のタイトルを変えて\n例：箇条書きをもっと短くして\n例：スライドを5枚追加して",
            height=100,
            key="slide_edit_instruction",
        )

        if st.button("スライドを修正する", use_container_width=True):
            if not slide_edit.strip():
                st.warning("修正指示を入力してください")
            else:
                try:
                    client = anthropic.Anthropic(api_key=api_key)
                    current_json = json.dumps(st.session_state.slide_data, ensure_ascii=False, indent=2)
                    revise_prompt = f"""以下のスライドデータを修正指示に従って修正してください。

## 現在のスライドデータ（JSON）
{current_json}

## 修正指示
{slide_edit}

## 出力ルール
- JSONのみ出力。前後の説明文・コードブロック記号は不要。
- 元のJSON構造を維持してください。"""

                    with st.spinner("スライドを修正中..."):
                        response = client.messages.create(
                            model=MODEL,
                            max_tokens=32000,
                            messages=[{"role": "user", "content": revise_prompt}]
                        )
                        text = response.content[0].text
                        slide_data = try_parse_json(text)
                        st.session_state.slide_data = slide_data

                    slide_count = len(slide_data.get("slides", []))
                    st.success(f"修正完了：{slide_count} 枚")
                    design = st.session_state.get("slide_design", DESIGN_PRESETS["ダーク（黒背景・白文字）"])
                    fmt = SLIDE_FORMATS.get(st.session_state.get("slide_format", list(SLIDE_FORMATS.keys())[0]))
                    rev_ofmts = st.session_state.get("slide_output_formats", ["PPT (.pptx)"])
                    display_name = st.session_state.display_name
                    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
                    # 修正後ZIPを再生成
                    with st.spinner("出力ZIPを生成中..."):
                        rev_zip = build_slides_zip(slide_data, design, fmt, rev_ofmts, display_name)
                    st.session_state.slide_zip_bytes = rev_zip
                    st.session_state.slide_approved = True
                    st.session_state.pop("slide_output_ts", None)

                    st.download_button(
                        "修正済みスライドをダウンロード (.zip)",
                        data=rev_zip,
                        file_name=f"{ts}_{display_name}_revised_slides.zip",
                        mime="application/zip",
                        use_container_width=True,
                        key="dl_revised_zip",
                    )

                except Exception as e:
                    st.error(f"スライド修正エラー: {e}")
